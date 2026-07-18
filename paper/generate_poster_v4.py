"""
CRAFT Poster v4  –  REaR-style
• Navy/blue palette + amber gold
• All 4 affiliation logos in header
• ACL 2026 banner embedded
• Large fonts (≥14 pt body), minimal text, maximum visuals
• Zigzag pipeline with coloured stage blocks + bold output badges
• Clean tables, big callout numbers, stat boxes
36" × 48" portrait
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x12, 0x2B, 0x4E)
BLUE       = RGBColor(0x1A, 0x5A, 0xA8)
BLUE_MED   = RGBColor(0x2E, 0x70, 0xC0)
BLUE_LT    = RGBColor(0xD6, 0xE8, 0xFA)
TEAL       = RGBColor(0x00, 0x7A, 0x87)
GOLD       = RGBColor(0xF5, 0xA6, 0x23)
GOLD_LT    = RGBColor(0xFF, 0xF0, 0xC0)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE  = RGBColor(0xF5, 0xF7, 0xFA)
LT_GRAY    = RGBColor(0xE2, 0xE8, 0xEF)
MID_GRAY   = RGBColor(0xB0, 0xBC, 0xCC)
DARK_TXT   = RGBColor(0x12, 0x1E, 0x2E)
GREEN_DK   = RGBColor(0x14, 0x65, 0x32)
GREEN_LT   = RGBColor(0xE6, 0xF5, 0xEB)
PURPLE     = RGBColor(0x5E, 0x17, 0x8E)
PURPLE_LT  = RGBColor(0xF0, 0xE6, 0xFA)
ORANGE     = RGBColor(0xC8, 0x56, 0x00)
ORANGE_LT  = RGBColor(0xFF, 0xF3, 0xE0)
RED_DK     = RGBColor(0xAA, 0x11, 0x11)
# Stage colours
P_BG, P_BD = RGBColor(0xFF,0xF8,0xEB), RGBColor(0xC8,0x78,0x00)   # pre-proc
S1_BG,S1_BD = GREEN_LT,  RGBColor(0x1B,0x7A,0x3E)
S2_BG,S2_BD = BLUE_LT,   RGBColor(0x1A,0x5A,0xA8)
S3_BG,S3_BD = PURPLE_LT, RGBColor(0x6A,0x1B,0x9A)
AG_BG,AG_BD = RGBColor(0xE8,0xF8,0xF1), RGBColor(0x00,0x65,0x4F)

# ── Dimensions ────────────────────────────────────────────────────────────
IN  = 914400
W   = int(36 * IN);  H  = int(48 * IN)
M   = int(0.22 * IN)
HDR = int(5.00 * IN)          # taller header to fit logos
BY  = HDR + int(0.16 * IN)
BH  = H - BY - M
GAP = int(0.20 * IN)
CW  = (W - 2*M - 2*GAP) // 3
CX  = [M + i*(CW+GAP) for i in range(3)]
SHH = int(0.62 * IN)          # section header height
IMGS = "/mnt/data1/asing725/ACL/CRAFT/static/images/"


# ── Primitives ────────────────────────────────────────────────────────────

def R(slide, x, y, w, h, fill=None, bd=None, bw=Pt(0.5)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.line.width = bw
    if fill is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb = fill
    if bd: s.line.color.rgb = bd
    else:  s.line.fill.background()
    return s

def T(slide, x, y, w, h, text, sz=Pt(15), bold=False, color=DARK_TXT,
      align=PP_ALIGN.LEFT, italic=False, wrap=True):
    t = slide.shapes.add_textbox(x, y, w, h)
    tf = t.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = sz; r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return t

def addP(tf, text, sz=Pt(15), bold=False, color=DARK_TXT,
         align=PP_ALIGN.LEFT, italic=False, sa=Pt(2)):
    p = tf.add_paragraph(); p.alignment = align; p.space_after = sa
    r = p.add_run(); r.text = text
    r.font.size = sz; r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color

def hdr(slide, x, y, w, title, fill=BLUE, tc=WHITE, h=None):
    hh = h or SHH
    R(slide, x, y, w, hh, fill=fill, bw=Pt(0))
    T(slide, x+int(0.14*IN), y+int(0.05*IN), w-int(0.18*IN), hh-int(0.08*IN),
      title, sz=Pt(19), bold=True, color=tc)
    return hh, y+hh

def pic(slide, path, x, y, w, h=None):
    try:
        if h: return slide.shapes.add_picture(path, x, y, w, h)
        else: return slide.shapes.add_picture(path, x, y, w)
    except Exception as e:
        T(slide, x, y, w, int(0.4*IN), f"[{path.split('/')[-1]}]", sz=Pt(10), color=MID_GRAY)


# ── HEADER ────────────────────────────────────────────────────────────────

def build_header(slide):
    # Navy background
    R(slide, 0, 0, W, HDR, fill=NAVY, bw=Pt(0))
    # Gold accent stripe
    R(slide, 0, HDR-int(0.09*IN), W, int(0.09*IN), fill=GOLD, bw=Pt(0))

    # ACL banner (right side)
    pic(slide, IMGS+"acl2026_banner.png",
        W - int(8.6*IN), int(0.18*IN), int(8.3*IN), int(1.55*IN))

    # Title (left, large)
    t = slide.shapes.add_textbox(int(0.28*IN), int(0.20*IN),
                                  int(27.5*IN), int(1.95*IN))
    tf = t.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "CRAFT: Training-Free Cascaded Retrieval for Tabular QA"
    r.font.size = Pt(50); r.font.bold = True; r.font.color.rgb = WHITE

    # Authors
    T(slide, int(0.28*IN), int(2.22*IN), int(27.8*IN), int(0.70*IN),
      "  *Adarsh Singh¹    *Kushal Raj Bhandari²    Jianxi Gao³    Soham Dan²    †Vivek Gupta¹",
      sz=Pt(20), bold=False, color=GOLD)

    # Affiliation line
    T(slide, int(0.28*IN), int(2.98*IN), int(27.8*IN), int(0.45*IN),
      "  ¹Arizona State University    ²Rensselaer Polytechnic Institute    "
      "³Microsoft    *Equal contribution    †Primary supervisor",
      sz=Pt(14.5), color=RGBColor(0xBB,0xCC,0xDD))

    # ── 4 Affiliation logos (side by side, white pill bg each) ───────────────
    # (path, display_width, display_height)
    logo_specs = [
        (IMGS+"asu_logo_hq.png",    int(3.60*IN), int(1.05*IN)),
        (IMGS+"rpi_logo.png",       int(2.60*IN), int(1.00*IN)),
        (IMGS+"microsoft_logo.png", int(3.40*IN), int(0.85*IN)),
        (IMGS+"coral_logo.jpeg",    int(1.10*IN), int(1.05*IN)),
    ]
    logo_y  = int(3.58*IN)
    pad_x   = int(0.10*IN); pad_y = int(0.08*IN)
    lx = int(0.28*IN)
    for path, lw, lh in logo_specs:
        pill_h = lh + 2*pad_y
        # White background pill
        R(slide, lx - pad_x, logo_y - pad_y,
          lw + 2*pad_x, pill_h, fill=WHITE, bw=Pt(0))
        # Logo centred vertically in pill
        pic(slide, path, lx, logo_y + (pill_h - lh)//2 - pad_y, lw, lh)
        lx += lw + int(0.42*IN)

    # arXiv box
    ax = W - int(5.8*IN); ay = int(1.90*IN)
    R(slide, ax, ay, int(5.5*IN), int(0.60*IN), fill=BLUE_MED, bw=Pt(0))
    T(slide, ax+int(0.12*IN), ay+int(0.05*IN), int(5.3*IN), int(0.55*IN),
      "arXiv: 2505.14984  ·  ACL 2026  ·  Vienna",
      sz=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Emails
    T(slide, ax+int(0.08*IN), ay+int(0.68*IN), int(5.4*IN), int(0.45*IN),
      "asing725@asu.edu   ·   bhandk@rpi.edu",
      sz=Pt(13.5), color=RGBColor(0xAA,0xBB,0xCC), align=PP_ALIGN.CENTER)


# ── MOTIVATION ────────────────────────────────────────────────────────────

def sec_motivation(slide, x, y, w, h):
    hh, by = hdr(slide, x, y, w, "  Motivation")
    R(slide, x, by, w, h-hh, fill=OFF_WHITE, bd=MID_GRAY, bw=Pt(0.5))
    bx = x+int(0.14*IN); bw = w-int(0.28*IN)

    t = slide.shapes.add_textbox(bx, by+int(0.12*IN), bw, h-hh-int(0.16*IN))
    tf = t.text_frame; tf.word_wrap = True

    def P(txt, bold=False, sz=Pt(15), c=DARK_TXT, sa=Pt(3)):
        addP(tf, txt, sz=sz, bold=bold, color=c, sa=sa)

    p = tf.paragraphs[0]; p.space_after = Pt(2)
    r = p.add_run(); r.text = "Open-Domain Table QA"
    r.font.size = Pt(16.5); r.font.bold = True; r.font.color.rgb = NAVY

    P("Retrieve the right table from 169K–419K tables, then answer natural language queries.", sz=Pt(14.5))
    P("", sa=Pt(4))
    P("Motivating Example", bold=True, sz=Pt(16.5), c=NAVY)
    P('Query: "Find Customers who bought products from suppliers in Japan"', sz=Pt(14.5), c=TEAL)
    P("  Top-2 retrieved:  Products,  Suppliers  →  No join path  ✗  SQL fails!", sz=Pt(14.5))
    P("  Missing:  Customers  ✓", sz=Pt(14.5), c=GREEN_DK)
    P("", sa=Pt(4))
    P("The Gap", bold=True, sz=Pt(16.5), c=NAVY)
    P("▸  Sparse/Dense retrieval ignores table–table joinability", sz=Pt(14.5))
    P("▸  LLM-based methods (ARM, JAR) are expensive at inference", sz=Pt(14.5))
    P("▸  Larger top-k hurts precision with noisy candidates", sz=Pt(14.5))
    P("", sa=Pt(4))
    P("Can pretrained off-the-shelf models, in a zero-shot cascade,\nmatch fine-tuned retrievers?",
      bold=True, sz=Pt(15), c=BLUE_MED)


# ── WHY EXISTING FAILS ────────────────────────────────────────────────────

def sec_fails(slide, x, y, w, h):
    hh, by = hdr(slide, x, y, w, "  Why Existing Approaches Fail")
    R(slide, x, by, w, h-hh, fill=OFF_WHITE, bd=MID_GRAY, bw=Pt(0.5))
    bx = x+int(0.10*IN); bw = w-int(0.20*IN)

    hdrs  = ["Approach", "Problem"]
    rows  = [
        ("Dense / Sparse",     "Optimises query–table relevance;\nignores joinability & connectivity"),
        ("LLM-based\n(ARM, JAR)", "Expensive inference loops →\nhigh latency & token cost"),
        ("Larger Top-k",       "Boosts recall but noisy tables\nhurt downstream precision"),
    ]
    c1 = int(bw*0.36); c2 = bw - c1
    row_h = int((h - hh - int(0.12*IN)) / (len(rows)+1))
    ty = by + int(0.06*IN)

    # Header
    for ci, (txt, cw, cx_) in enumerate([(hdrs[0],c1,bx),(hdrs[1],c2,bx+c1)]):
        R(slide, cx_, ty, cw, row_h, fill=NAVY, bw=Pt(0))
        T(slide, cx_+int(0.07*IN), ty+int(0.04*IN), cw-int(0.1*IN), row_h-int(0.06*IN),
          txt, sz=Pt(15), bold=True, color=WHITE)
    ty += row_h

    row_fills = [LT_GRAY, BLUE_LT, ORANGE_LT]
    for ri, (a, b) in enumerate(rows):
        fill = row_fills[ri]
        for txt, cw, cx_ in [(a,c1,bx),(b,c2,bx+c1)]:
            R(slide, cx_, ty, cw, row_h, fill=fill, bd=MID_GRAY, bw=Pt(0.4))
            T(slide, cx_+int(0.07*IN), ty+int(0.04*IN), cw-int(0.10*IN), row_h-int(0.06*IN),
              txt, sz=Pt(14))
        ty += row_h


# ── KEY INSIGHT ───────────────────────────────────────────────────────────

def sec_insight(slide, x, y, w, h):
    hh, by = hdr(slide, x, y, w, "  Key Insight", fill=RGBColor(0xA0,0x6A,0x00), tc=WHITE)
    R(slide, x, by, w, h-hh, fill=GOLD_LT, bd=GOLD, bw=Pt(1.5))
    bx = x+int(0.14*IN); bw = w-int(0.28*IN)

    t = slide.shapes.add_textbox(bx, by+int(0.12*IN), bw, h-hh-int(0.16*IN))
    tf = t.text_frame; tf.word_wrap = True

    p0 = tf.paragraphs[0]; p0.space_after = Pt(3)
    r0 = p0.add_run(); r0.text = "CRAFT jointly models:"
    r0.font.size = Pt(16); r0.font.bold = True; r0.font.color.rgb = NAVY

    items = [
        ("Query ↔ Table  Relevance", "(What is semantically relevant?)", BLUE),
        ("Table ↔ Table  Joinability", "(Can they be joined together?)",  GREEN_DK),
        ("⟹  Join-ready, high-fidelity sets", "Without any LLM calls",       NAVY),
    ]
    for title, sub, c in items:
        addP(tf, title, sz=Pt(15.5), bold=True, color=c, sa=Pt(1))
        addP(tf, f"  {sub}", sz=Pt(14), color=DARK_TXT, sa=Pt(4))


# ── CONTRIBUTIONS ─────────────────────────────────────────────────────────

def sec_contributions(slide, x, y, w, h):
    hh, by = hdr(slide, x, y, w, "  Key Contributions")
    R(slide, x, by, w, h-hh, fill=OFF_WHITE, bd=MID_GRAY, bw=Pt(0.5))
    bx = x+int(0.14*IN); bw = w-int(0.28*IN)

    t = slide.shapes.add_textbox(bx, by+int(0.14*IN), bw, h-hh-int(0.18*IN))
    tf = t.text_frame; tf.word_wrap = True

    items = [
        ("① Zero-Shot · No Fine-Tuning",
         "Plug-and-play on any domain without dataset-specific training."),
        ("② SOTA on NQ-Tables",
         "R@1=49.84  R@10=86.83  R@50=97.17 — all zero-shot."),
        ("③ Robust to Paraphrasing",
         "Only –0.04 Δ on perturbed queries vs. –5–12 for DTR."),
        ("④ 33× Fewer Embedding Calls",
         "Mini-tables cut >70% tokens; no online LLM needed."),
    ]
    p0 = tf.paragraphs[0]
    r0 = p0.add_run(); r0.text = items[0][0]
    r0.font.size = Pt(16); r0.font.bold = True; r0.font.color.rgb = NAVY
    addP(tf, items[0][1], sz=Pt(14.5), sa=Pt(6))
    for title, body in items[1:]:
        addP(tf, title, sz=Pt(16), bold=True, color=NAVY, sa=Pt(1))
        addP(tf, body, sz=Pt(14.5), sa=Pt(6))


# ── OVERVIEW FIGURE ───────────────────────────────────────────────────────

def sec_overview(slide, x, y, w, h):
    hh, by = hdr(slide, x, y, w, "  CRAFT Framework Overview")
    R(slide, x, by, w, h-hh, fill=WHITE, bd=MID_GRAY, bw=Pt(0.6))
    pad = int(0.12*IN)
    pic(slide, IMGS+"craft_overview.png", x+pad, by+pad, w-2*pad, h-hh-2*pad)


# ── PIPELINE (ZIGZAG) ─────────────────────────────────────────────────────

def sec_pipeline(slide, x, y, w, h):
    hh, by = hdr(slide, x, y, w, "  CRAFT Pipeline  ·  Three-Stage Cascade")
    R(slide, x, by, w, h-hh, fill=OFF_WHITE, bd=MID_GRAY, bw=Pt(0.5))

    bx = x+int(0.12*IN); bw = w-int(0.24*IN)
    cy = by + int(0.14*IN)
    ARR = int(0.30*IN)   # arrow height

    def arrow(cy_, color):
        mid = bx + bw//2
        R(slide, mid-int(0.05*IN), cy_, int(0.10*IN), ARR-int(0.12*IN), fill=color, bw=Pt(0))
        T(slide, mid-int(0.18*IN), cy_+ARR-int(0.22*IN), int(0.36*IN), int(0.24*IN),
          "▼", sz=Pt(12), bold=True, color=color, align=PP_ALIGN.CENTER)
        return cy_ + ARR

    # ── Pre-processing ───────────────────────────────────────────────────
    ph = int(2.10*IN); badge_w = int(0.72*IN)
    R(slide, bx, cy, bw, ph, fill=P_BG, bd=P_BD, bw=Pt(1.5))
    R(slide, bx, cy, badge_w, ph, fill=P_BD, bw=Pt(0))
    T(slide, bx+int(0.03*IN), cy+int(0.08*IN), badge_w-int(0.04*IN), ph-int(0.10*IN),
      "PRE\nPROC", sz=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    t = slide.shapes.add_textbox(bx+badge_w+int(0.10*IN), cy+int(0.08*IN),
                                  bw-badge_w-int(0.14*IN), ph-int(0.12*IN))
    tf = t.text_frame; tf.word_wrap = True
    p0 = tf.paragraphs[0]; p0.space_after = Pt(3)
    r0 = p0.add_run(); r0.text = "Gemini Flash 1.5  (offline, once)"
    r0.font.size = Pt(15); r0.font.bold = True; r0.font.color.rgb = ORANGE
    for line in ["▸  Query Decomposition — break into sub-questions",
                 "▸  Table Enrichment — generate titles + summaries",
                 "▸  Row Ranking — rank rows by semantic relevance"]:
        addP(tf, line, sz=Pt(14), sa=Pt(2))
    cy += ph

    # ── 3 stages ─────────────────────────────────────────────────────────
    stages = [
        ("STAGE\n  1", "Sparse Lexical Retrieval  ·  SPLADE",
         ["▸  Full corpus:  169K (NQ) / 419K (OTT-QA) tables",
          "▸  Enriched titles + sub-questions boost recall",
          "▸  Inverted-index lookup — near zero latency"],
         "→ N₁ = 5,000\n  candidates", S1_BG, S1_BD),
        ("STAGE\n  2", "Dense Semantic Pruning  ·  Sentence Transformer",
         ["▸  all-mpnet-base-v2 (NQ)  /  Jina v3 (OTT-QA)",
          "▸  Mini-tables (top-5 rows) scored vs. query",
          "▸  Sublinear ANN search — pre-encoded offline"],
         "→ N₂ = 100–300\n  mini-tables", S2_BG, S2_BD),
        ("STAGE\n  3", "Neural Re-ranking  ·  Precision-Optimised",
         ["▸  text-embedding-3-small / large  (NQ-Tables)",
          "▸  gemini-embedding-001  (OTT-QA)",
          "▸  Plug-and-play: swap model with zero code changes"],
         "→ Top-K tables\n  K ∈ {1,3,5,8,10}", S3_BG, S3_BD),
    ]
    sh = int(2.55*IN); ob_w = int(1.75*IN)
    stage_bw = bw - ob_w - int(0.10*IN)

    for label, title, bullets, out_txt, bg, bd in stages:
        cy = arrow(cy, bd)
        # Stage block
        R(slide, bx, cy, stage_bw, sh, fill=bg, bd=bd, bw=Pt(1.8))
        R(slide, bx, cy, badge_w, sh, fill=bd, bw=Pt(0))
        T(slide, bx+int(0.03*IN), cy+int(0.08*IN), badge_w-int(0.04*IN), sh-int(0.10*IN),
          label, sz=Pt(15), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Title + bullets
        T(slide, bx+badge_w+int(0.10*IN), cy+int(0.06*IN),
          stage_bw-badge_w-int(0.14*IN), int(0.46*IN),
          title, sz=Pt(14.5), bold=True, color=bd)
        t2 = slide.shapes.add_textbox(bx+badge_w+int(0.10*IN), cy+int(0.50*IN),
                                       stage_bw-badge_w-int(0.14*IN), sh-int(0.55*IN))
        tf2 = t2.text_frame; tf2.word_wrap = True
        p1 = tf2.paragraphs[0]; p1.space_after = Pt(2)
        r1 = p1.add_run(); r1.text = bullets[0]
        r1.font.size = Pt(13.5); r1.font.color.rgb = DARK_TXT
        for b in bullets[1:]:
            addP(tf2, b, sz=Pt(13.5), sa=Pt(2))
        # Output badge
        ob_x = bx + stage_bw + int(0.10*IN)
        ob_y = cy + (sh - int(1.05*IN))//2
        R(slide, ob_x, ob_y, ob_w, int(1.05*IN), fill=bd, bw=Pt(0))
        T(slide, ob_x+int(0.07*IN), ob_y+int(0.05*IN),
          ob_w-int(0.12*IN), int(0.98*IN),
          out_txt, sz=Pt(13.5), bold=True, color=WHITE)
        cy += sh

    # Answer generation bar
    cy = arrow(cy, AG_BD)
    ag_h = int(1.18*IN)
    R(slide, bx, cy, bw, ag_h, fill=AG_BG, bd=AG_BD, bw=Pt(1.5))
    R(slide, bx, cy, badge_w, ag_h, fill=AG_BD, bw=Pt(0))
    T(slide, bx+int(0.03*IN), cy+int(0.04*IN), badge_w-int(0.04*IN), ag_h-int(0.06*IN),
      "ANS\nGEN", sz=Pt(13.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    T(slide, bx+badge_w+int(0.12*IN), cy+int(0.10*IN),
      bw-badge_w-int(0.20*IN), ag_h-int(0.14*IN),
      "Pretrained LLM  (Llama3-8B / Qwen2.5-7B / Mistral-7B)\n"
      "Top-K mini-tables  +  Query  →  Final Answer",
      sz=Pt(14.5), color=DARK_TXT)


# ── TABLE BUILDER ─────────────────────────────────────────────────────────

CAT_CLR = {
    "cat":    RGBColor(0xC2,0xD5,0xEE),
    "sparse": OFF_WHITE,
    "dense":  BLUE_LT,
    "hybrid": GREEN_LT,
    "craft_s":GOLD_LT,
    "craft":  RGBColor(0xFF,0xE5,0x80),
    "base":   LT_GRAY,
    "bad":    RGBColor(0xFF,0xE0,0xE0),
}

def build_table(slide, bx, ty, bw, fracs, headers, rows,
                rh=int(0.31*IN), hdr_fill=NAVY):
    cws = [int(bw*f) for f in fracs]

    def cell(cx_, ry, cw, ch, txt, fill, tc, bold, sz, al):
        R(slide, cx_, ry, cw, ch, fill=fill, bd=MID_GRAY, bw=Pt(0.4))
        T(slide, cx_+int(0.05*IN), ry+int(0.02*IN), cw-int(0.07*IN), ch-int(0.03*IN),
          txt, sz=sz, bold=bold, color=tc, align=al)

    cx_ = bx
    for ci, (h2, cw) in enumerate(zip(headers, cws)):
        cell(cx_, ty, cw, rh, h2, hdr_fill, WHITE, True, Pt(13.5),
             PP_ALIGN.CENTER if ci>0 else PP_ALIGN.LEFT)
        cx_ += cw
    ty += rh

    for row in rows:
        *cells, cat, bold = row
        fill = CAT_CLR.get(cat, LT_GRAY)
        tc = NAVY if (cat in ("cat","craft")) else DARK_TXT
        cx_ = bx
        for ci, (val, cw) in enumerate(zip(cells, cws)):
            cell(cx_, ty, cw, rh, val, fill, tc, bold,
                 Pt(14 if cat!="cat" else 12.5),
                 PP_ALIGN.CENTER if ci>0 else PP_ALIGN.LEFT)
            cx_ += cw
        ty += rh
    return ty


# ── NQ-TABLES RETRIEVAL ───────────────────────────────────────────────────

def sec_nq(slide, x, y, w, h):
    hh, by = hdr(slide, x, y, w, "  Retrieval  ·  NQ-Tables (169K tables, trained models excluded)")
    R(slide, x, by, w, h-hh, fill=OFF_WHITE, bd=MID_GRAY, bw=Pt(0.5))
    bx = x+int(0.10*IN); bw = w-int(0.20*IN)
    T(slide, bx, by+int(0.09*IN), bw, int(0.50*IN),
      "R = Recall  ·  Bold = best  ·  All non-CRAFT models trained on NQ-Tables",
      sz=Pt(13.5), italic=True, color=DARK_TXT)
    rows = [
        ("Sparse Retrieval", "", "", "", "cat", True),
        ("BM25",          "18.49","36.94","52.61","sparse",False),
        ("SPLADE",        "39.84","83.33","94.65","sparse",False),
        ("Dense Retrieval","","","","cat",True),
        ("BIBERT",        "43.78","82.25","93.71","dense",False),
        ("DPR",           "45.32","85.84","95.44","dense",False),
        ("DTR",           "32.62","75.86","89.77","dense",False),
        ("SSDR",          "45.47","84.00","95.05","dense",False),
        ("Hybrid Retrieval","","","","cat",True),
        ("BIBERT+SPLADE", "43.67","86.72","95.62","hybrid",False),
        ("THYME",         "48.55","86.38","96.08","hybrid",False),
        ("CRAFT Stages",  "","","","cat",True),
        ("Stage 1",       "34.38","72.90","91.62","craft_s",False),
        ("+ Stage 2",     "36.65","82.91","96.08","craft_s",False),
        ("+ Stage 3",     "41.13","87.16","96.84","craft_s",False),
        ("CRAFT (full) ★","49.84","86.83","97.17","craft",True),
    ]
    rh = min(int((h-hh-int(0.60*IN)-int(0.70*IN))/len(rows)), int(0.30*IN))
    ty = build_table(slide, bx, by+int(0.62*IN), bw,
                     [0.43,0.19,0.19,0.19], ["Model","R@1","R@10","R@50"],
                     rows, rh=rh)
    rem = (y+h) - ty - int(0.06*IN)
    if rem > int(0.38*IN):
        ch = min(int(0.60*IN), rem)
        R(slide, bx, ty+int(0.06*IN), bw, ch, fill=GOLD_LT, bd=GOLD, bw=Pt(1.5))
        T(slide, bx+int(0.10*IN), ty+int(0.10*IN), bw-int(0.15*IN), ch-int(0.10*IN),
          "★  CRAFT  R@1 = 49.84  ·  new SOTA  ·  no training required",
          sz=Pt(14.5), bold=True, color=NAVY)


# ── OTT-QA RETRIEVAL ─────────────────────────────────────────────────────

def sec_ottqa(slide, x, y, w, h):
    hh, by = hdr(slide, x, y, w, "  Retrieval  ·  OTT-QA  (419K tables · Zero-Shot)")
    R(slide, x, by, w, h-hh, fill=OFF_WHITE, bd=MID_GRAY, bw=Pt(0.5))
    bx = x+int(0.10*IN); bw = w-int(0.20*IN)
    T(slide, bx, by+int(0.09*IN), bw, int(0.50*IN),
      "Multi-hop reasoning across tables + passages. CRAFT fully zero-shot.",
      sz=Pt(13.5), italic=True, color=DARK_TXT)
    rows = [
        ("Sparse Retrieval","","","","cat",True),
        ("BM25",           "23.98","51.94","69.11","sparse",False),
        ("SPLADE",         "62.74","89.52","95.21","sparse",False),
        ("Dense Retrieval","","","","cat",True),
        ("BIBERT",         "56.82","86.50","94.26","dense",False),
        ("DPR",            "53.43","85.95","93.22","dense",False),
        ("SSDR",           "56.96","86.22","93.55","dense",False),
        ("Hybrid Retrieval","","","","cat",True),
        ("BIBERT+SPLADE",  "64.72","91.01","96.34","hybrid",False),
        ("THYME",          "66.67","91.10","96.16","hybrid",False),
        ("CRAFT Stages",   "","","","cat",True),
        ("Stage 1",        "40.74","60.89","88.17","craft_s",False),
        ("+ Stage 2",      "47.33","87.35","91.82","craft_s",False),
        ("CRAFT (full)",   "55.56","89.88","96.07","craft",True),
    ]
    rh = min(int((h-hh-int(0.60*IN)-int(0.68*IN))/len(rows)), int(0.30*IN))
    ty = build_table(slide, bx, by+int(0.62*IN), bw,
                     [0.43,0.19,0.19,0.19], ["Model","R@1","R@10","R@50"],
                     rows, rh=rh)
    rem = (y+h) - ty - int(0.06*IN)
    if rem > int(0.36*IN):
        ch = min(int(0.58*IN), rem)
        R(slide, bx, ty+int(0.06*IN), bw, ch, fill=BLUE_LT, bd=BLUE, bw=Pt(1.2))
        T(slide, bx+int(0.10*IN), ty+int(0.10*IN), bw-int(0.15*IN), ch-int(0.10*IN),
          "Zero-shot CRAFT approaches fine-tuned THYME at R@10 & R@50",
          sz=Pt(14.5), bold=True, color=NAVY)


# ── END-TO-END QA ─────────────────────────────────────────────────────────

def sec_e2e(slide, x, y, w, h):
    hh, by = hdr(slide, x, y, w,
                 "  End-to-End QA Performance (F1)  ·  NQ-Tables")
    R(slide, x, by, w, h-hh, fill=OFF_WHITE, bd=MID_GRAY, bw=Pt(0.5))
    bx = x+int(0.10*IN); bw = w-int(0.20*IN)
    T(slide, bx, by+int(0.09*IN), bw, int(0.50*IN),
      "F1 · n = number of retrieved tables · Models: Mistral-7B, Llama3-8B, Qwen2.5-7B",
      sz=Pt(13.5), italic=True, color=DARK_TXT)

    rh = int(0.31*IN)
    cws_raw = [0.24,0.085,0.085,0.085,0.085,0.085,0.085,0.085,0.085,0.085]
    cws = [int(bw*f) for f in cws_raw]

    # Span header
    ty = by + int(0.62*IN)
    spans = [("", 1, NAVY), ("n = 1", 3, BLUE_MED),
             ("n = 3", 3, GREEN_DK), ("n = 5", 3, PURPLE)]
    cx_ = bx
    widths_used = 0
    for lbl, n, c in spans:
        sw = sum(cws[widths_used:widths_used+n])
        R(slide, cx_, ty, sw, rh, fill=c, bw=Pt(0))
        T(slide, cx_, ty, sw, rh, lbl, sz=Pt(13.5), bold=True,
          color=WHITE, align=PP_ALIGN.CENTER)
        cx_ += sw; widths_used += n
    ty += rh

    sub_h = ["Retriever","Mis.","Ll3","Qwn","Mis.","Ll3","Qwn","Mis.","Ll3","Qwn"]
    sub_c = [NAVY]+[BLUE_MED]*3+[GREEN_DK]*3+[PURPLE]*3
    cx_ = bx
    for sh, cw, sc in zip(sub_h, cws, sub_c):
        R(slide, cx_, ty, cw, rh, fill=RGBColor(0xCC,0xD8,0xEE), bd=MID_GRAY, bw=Pt(0.3))
        T(slide, cx_, ty, cw, rh, sh, sz=Pt(13), bold=True, color=sc, align=PP_ALIGN.CENTER)
        cx_ += cw
    ty += rh

    data = [
        ("BIBERT",       "32.93","32.66","34.80","34.6","33.16","37.92","35.30","34.28","37.01","base",False),
        ("SPLADE",       "29.61","32.07","31.90","35.42","37.17","35.79","34.95","33.88","37.35","base",False),
        ("BIBERT+SPLADE","32.67","32.66","33.24","33.59","33.66","36.76","35.71","33.92","37.02","base",False),
        ("THYME",        "35.48","36.14","37.28","37.59","39.16","40.28","37.20","39.29","41.20","hybrid",False),
        ("CRAFT ★",      "39.13","38.31","39.73","45.28","40.76","43.52","44.53","40.55","46.49","craft",True),
    ]
    row_f = {"base":LT_GRAY, "hybrid":GREEN_LT, "craft":RGBColor(0xFF,0xE5,0x80)}
    for row in data:
        *cells, cat, bold = row
        fill = row_f.get(cat, LT_GRAY)
        tc = NAVY if bold else DARK_TXT
        cx_ = bx
        for ci, (val, cw) in enumerate(zip(cells, cws)):
            R(slide, cx_, ty, cw, rh, fill=fill, bd=MID_GRAY, bw=Pt(0.3))
            T(slide, cx_+int(0.03*IN), ty+int(0.02*IN), cw-int(0.05*IN), rh-int(0.03*IN),
              val, sz=Pt(13.5 if ci==0 else 12.5), bold=bold, color=tc,
              align=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER)
            cx_ += cw
        ty += rh

    rem = (y+h) - ty - int(0.06*IN)
    if rem > int(0.38*IN):
        ch = min(int(0.58*IN), rem)
        R(slide, bx, ty+int(0.06*IN), bw, ch, fill=GOLD_LT, bd=GOLD, bw=Pt(1.5))
        T(slide, bx+int(0.10*IN), ty+int(0.10*IN), bw-int(0.15*IN), ch-int(0.10*IN),
          "★  CRAFT achieves best F1 at n=1 and n=5 across all three LLMs",
          sz=Pt(14.5), bold=True, color=NAVY)


# ── ABLATION BAR CHART ────────────────────────────────────────────────────

def sec_ablation(slide, x, y, w, h):
    hh, by = hdr(slide, x, y, w, "  Ablation  ·  Stage-wise R@10 Gains",
                 fill=PURPLE)
    R(slide, x, by, w, h-hh, fill=OFF_WHITE, bd=MID_GRAY, bw=Pt(0.5))
    bx = x+int(0.12*IN); bw = w-int(0.24*IN)

    T(slide, bx, by+int(0.10*IN), bw, int(0.44*IN),
      "NQ-Tables Recall@10 — each stage adds distinct gains:",
      sz=Pt(15), bold=True, color=DARK_TXT)

    bars = [("Stage 1\nSPLADE", 72.90, S1_BD),
            ("+Stage 2\nDense",  82.91, S2_BD),
            ("+Stage 3\nRerank", 87.16, S3_BD),
            ("CRAFT\n(full)",    86.83, NAVY)]

    chart_y = by + int(0.58*IN)
    chart_h = h - hh - int(0.60*IN) - int(0.55*IN)
    axis_h  = chart_h - int(0.48*IN)
    bpad    = int(0.35*IN)
    bar_w   = (bw - bpad - int(0.10*IN)) // len(bars) - int(0.12*IN)
    bx0     = bx + bpad

    for i, (label, val, c) in enumerate(bars):
        xi = bx0 + i*(bar_w+int(0.12*IN))
        fh = int((val/100.0)*axis_h)
        bt = chart_y + axis_h - fh
        # Track
        R(slide, xi, chart_y, bar_w, axis_h, fill=LT_GRAY, bw=Pt(0))
        # Bar
        R(slide, xi, bt, bar_w, fh, fill=c, bw=Pt(0))
        # Value
        T(slide, xi-int(0.06*IN), bt-int(0.35*IN), bar_w+int(0.12*IN), int(0.33*IN),
          f"{val:.1f}", sz=Pt(14.5), bold=True, color=c, align=PP_ALIGN.CENTER)
        # Label
        T(slide, xi-int(0.04*IN), chart_y+axis_h+int(0.04*IN),
          bar_w+int(0.08*IN), int(0.44*IN),
          label, sz=Pt(12.5), color=DARK_TXT, align=PP_ALIGN.CENTER)

    # Y-axis lines
    for val in [50, 75, 100]:
        fy = chart_y + axis_h - int((val/100.0)*axis_h)
        R(slide, bx, fy, bw, int(0.015*IN), fill=MID_GRAY, bw=Pt(0))
        T(slide, bx, fy-int(0.20*IN), bpad-int(0.05*IN), int(0.22*IN),
          str(val), sz=Pt(12), color=MID_GRAY, align=PP_ALIGN.RIGHT)

    insight_y = chart_y + chart_h + int(0.05*IN)
    rem = (y+h) - insight_y - int(0.05*IN)
    if rem > int(0.36*IN):
        ch = min(int(0.52*IN), rem)
        R(slide, bx, insight_y, bw, ch, fill=PURPLE_LT, bd=PURPLE, bw=Pt(1.0))
        T(slide, bx+int(0.10*IN), insight_y+int(0.06*IN), bw-int(0.14*IN), ch-int(0.08*IN),
          "Each stage delivers distinct, non-overlapping recall improvements",
          sz=Pt(14), bold=True, color=PURPLE)


# ── ROBUSTNESS ────────────────────────────────────────────────────────────

def sec_robustness(slide, x, y, w, h):
    hh, by = hdr(slide, x, y, w, "  Robustness to Query Paraphrasing",
                 fill=RGBColor(0x10,0x55,0x9A))
    R(slide, x, by, w, h-hh, fill=OFF_WHITE, bd=MID_GRAY, bw=Pt(0.5))
    bx = x+int(0.10*IN); bw = w-int(0.20*IN)
    T(slide, bx, by+int(0.09*IN), bw, int(0.50*IN),
      "Queries paraphrased by Gemini 2.5 Flash.  Δ = change in avg. Recall vs. original.",
      sz=Pt(13.5), italic=True, color=DARK_TXT)

    rows = [
        ("DTR (M) — original",  "38.74","75.73","88.36","—",  "base",False),
        ("CRAFT — original",    "41.13","87.16","96.84","—",  "craft",True),
        ("DTR (M) — perturbed", "27.09","66.27","84.33","–8.38","bad",False),
        ("DTR+HN — perturbed",  "37.87","76.06","84.43","–5.80","bad",False),
        ("CRAFT — perturbed",   "41.02","86.83","96.08","–0.04","craft",True),
    ]
    cat_r = {**CAT_CLR, "bad": RGBColor(0xFF,0xE0,0xE0), "base": BLUE_LT}
    rh = int(0.31*IN)
    ty = build_table(slide, bx, by+int(0.62*IN), bw,
                     [0.38,0.15,0.15,0.15,0.17],
                     ["Model","R@1","R@10","R@50","Δ Avg"],
                     rows, rh=rh, hdr_fill=RGBColor(0x10,0x55,0x9A))

    # Override cat_fills for bad
    rem = (y+h) - ty - int(0.06*IN)
    if rem > int(0.36*IN):
        ch = min(int(0.58*IN), rem)
        R(slide, bx, ty+int(0.06*IN), bw, ch, fill=BLUE_LT, bd=BLUE, bw=Pt(1.2))
        T(slide, bx+int(0.10*IN), ty+int(0.10*IN), bw-int(0.15*IN), ch-int(0.10*IN),
          "CRAFT: only –0.04 Δ on paraphrased queries  vs.  –5.8 to –8.4 for DTR",
          sz=Pt(14.5), bold=True, color=NAVY)


# ── EFFICIENCY ────────────────────────────────────────────────────────────

def sec_efficiency(slide, x, y, w, h):
    hh, by = hdr(slide, x, y, w, "  Efficiency  ·  Mini-Tables & Embedding Calls",
                 fill=GREEN_DK)
    R(slide, x, by, w, h-hh, fill=OFF_WHITE, bd=MID_GRAY, bw=Pt(0.5))
    bx = x+int(0.12*IN); bw = w-int(0.24*IN)

    # 3 big stat boxes
    stats = [("33×", "fewer\nembedding\ncalls",     S1_BD),
             (">70%","token\nreduction",            BLUE_MED),
             ("0",   "online\nLLM\ncalls",          PURPLE)]
    sw = (bw - int(0.16*IN)) // 3
    sx = bx; sy = by + int(0.12*IN)
    for big, small, c in stats:
        R(slide, sx, sy, sw, int(1.38*IN), fill=WHITE, bd=c, bw=Pt(2.0))
        T(slide, sx, sy+int(0.04*IN), sw, int(0.72*IN),
          big, sz=Pt(32), bold=True, color=c, align=PP_ALIGN.CENTER)
        T(slide, sx+int(0.04*IN), sy+int(0.72*IN), sw-int(0.06*IN), int(0.62*IN),
          small, sz=Pt(13), color=DARK_TXT, align=PP_ALIGN.CENTER)
        sx += sw + int(0.08*IN)

    t = slide.shapes.add_textbox(bx, by+int(1.60*IN), bw, h-hh-int(1.65*IN))
    tf = t.text_frame; tf.word_wrap = True

    def P(txt, bold=False, sz=Pt(14.5), c=DARK_TXT, sa=Pt(3)):
        addP(tf, txt, sz=sz, bold=bold, color=c, sa=sa)

    p0 = tf.paragraphs[0]; p0.space_after = Pt(2)
    r0 = p0.add_run(); r0.text = "Token Reduction via Mini-Tables"
    r0.font.size = Pt(16); r0.font.bold = True; r0.font.color.rgb = NAVY
    P("Mini-tables (top-5 rows) cut context by >70%:")
    P("  Llama3: 1,783 → 363 tokens  ·  Mistral: 3,314 → 477 tokens")
    P("", sa=Pt(3))
    P("Embedding Call Reduction", bold=True, sz=Pt(16), c=NAVY)
    P("Dense baseline: 169,898 table embeddings per query")
    P("CRAFT: 3 query embeddings + 5,100 mini-table embeddings")
    P("33× leaner — bulk of computation is offline & pre-computed.", bold=True, c=GREEN_DK)
    P("", sa=Pt(3))
    P("Adaptability", bold=True, sz=Pt(16), c=NAVY)
    P("Stage 3 is plug-and-play — swap in any new embedding model\n"
      "with zero architectural changes.")


# ── CONCLUSION ────────────────────────────────────────────────────────────

def sec_conclusion(slide, x, y, w, h):
    hh, by = hdr(slide, x, y, w, "  Conclusion", fill=NAVY)
    R(slide, x, by, w, h-hh, fill=BLUE_LT, bd=BLUE, bw=Pt(0.8))
    bx = x+int(0.14*IN); bw = w-int(0.28*IN)

    t = slide.shapes.add_textbox(bx, by+int(0.12*IN), bw, h-hh-int(0.16*IN))
    tf = t.text_frame; tf.word_wrap = True

    def P(txt, bold=False, sz=Pt(14.5), c=DARK_TXT, sa=Pt(3)):
        addP(tf, txt, sz=sz, bold=bold, color=c, sa=sa)

    p0 = tf.paragraphs[0]; p0.space_after = Pt(3)
    r0 = p0.add_run()
    r0.text = ("CRAFT is a modular, training-free cascade for Open-Domain "
               "Table QA that achieves SOTA recall without fine-tuning.")
    r0.font.size = Pt(15); r0.font.color.rgb = DARK_TXT
    P("", sa=Pt(4))
    P("Key Takeaways", bold=True, sz=Pt(16.5), c=NAVY)
    P("▸  R@1 = 49.84 — new SOTA on NQ-Tables",    sz=Pt(15))
    P("▸  Competitive zero-shot on OTT-QA",          sz=Pt(15))
    P("▸  33× fewer embedding calls",                sz=Pt(15))
    P("▸  –0.04 Δ under query paraphrasing",         sz=Pt(15))
    P("▸  Plug-and-play: swap any Stage-3 model",    sz=Pt(15))
    P("", sa=Pt(4))
    P("Bridges fine-tuned & lightweight retrieval —\n"
      "scalable, interpretable, adaptable.",
      bold=True, sz=Pt(15.5), c=BLUE_MED)


# ── MAIN ──────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = Emu(W)
    prs.slide_height = Emu(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full background
    R(slide, 0, 0, W, H, fill=RGBColor(0xEC, 0xF1, 0xF8))

    build_header(slide)

    # ── Body layout (BODY_Y → H-M) ────────────────────────────────────────
    #  Band A (tall):  left col  |  center col  |  right col
    #  Band B:         e2e (2 cols)              |  ablation
    #  Band C:         robustness | efficiency   | conclusion

    A_H = int(27.2 * IN)
    B_H = int( 8.6 * IN)
    C_H = BH - A_H - B_H - 2*GAP

    AY = BY;              BY2 = AY + A_H + GAP;  CY = BY2 + B_H + GAP

    # ── Left column (Band A) ──────────────────────────────────────────────
    sections_left = [
        (int(10.5*IN), sec_motivation),
        (int( 5.8*IN), sec_fails),
        (int( 4.8*IN), sec_insight),
    ]
    remaining_a = A_H - sum(s for s,_ in sections_left) - (len(sections_left)-1)*GAP
    sections_left.append((remaining_a, sec_contributions))

    oy = AY
    for sh, fn in sections_left:
        fn(slide, CX[0], oy, CW, sh)
        oy += sh + GAP

    # ── Center column (Band A): overview figure (top) + pipeline ─────────
    FIG_H = int(10.6 * IN)
    PIP_H = A_H - FIG_H - GAP
    sec_overview(slide, CX[1], AY, CW, FIG_H)
    sec_pipeline(slide, CX[1], AY+FIG_H+GAP, CW, PIP_H)

    # ── Right column (Band A): NQ + OTT ───────────────────────────────────
    NQ_H  = int(14.0 * IN)
    OTT_H = A_H - NQ_H - GAP
    sec_nq(   slide, CX[2], AY,         CW, NQ_H)
    sec_ottqa(slide, CX[2], AY+NQ_H+GAP, CW, OTT_H)

    # ── Band B ────────────────────────────────────────────────────────────
    e2e_w = CW*2 + GAP
    sec_e2e(    slide, CX[0], BY2, e2e_w, B_H)
    sec_ablation(slide, CX[2], BY2, CW,    B_H)

    # ── Band C ────────────────────────────────────────────────────────────
    sec_robustness(slide, CX[0], CY, CW, C_H)
    sec_efficiency(slide, CX[1], CY, CW, C_H)
    sec_conclusion(slide, CX[2], CY, CW, C_H)

    # ── Footer ────────────────────────────────────────────────────────────
    FY = H - int(0.32*IN)
    R(slide, 0, FY, W, int(0.32*IN), fill=NAVY, bw=Pt(0))
    T(slide, int(0.3*IN), FY+int(0.04*IN), int(22*IN), int(0.26*IN),
      "CRAFT  ·  Training-Free Cascaded Retrieval for Tabular QA  ·  ACL 2026",
      sz=Pt(12.5), color=RGBColor(0xBB,0xCC,0xDD))
    T(slide, int(14*IN), FY+int(0.04*IN), int(21.5*IN), int(0.26*IN),
      "arXiv: 2505.14984   ·   asing725@asu.edu   ·   bhandk@rpi.edu   ·   "
      "coral-lab-asu.github.io",
      sz=Pt(12.5), color=GOLD, align=PP_ALIGN.RIGHT)

    out = "/mnt/data1/asing725/ACL/CRAFT/CRAFT_poster.pptx"
    prs.save(out)
    ns = len(slide.shapes)
    np_ = sum(1 for s in slide.shapes if s.shape_type == 13)
    print(f"Saved  {out}")
    print(f"Size   {prs.slide_width.inches:.0f}\" × {prs.slide_height.inches:.0f}\"  "
          f"|  {ns} shapes  |  {np_} images")


if __name__ == "__main__":
    build()
