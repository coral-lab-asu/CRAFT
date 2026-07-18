"""
CRAFT Conference Poster Generator (v2)
Inspired by the Weaver poster style: colored section headers, clean tables,
modular layout with strong visual hierarchy.
Poster size: 36" x 48" (portrait) — standard ACL poster
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy
from lxml import etree

# ── Color palette (ASU Maroon + Gold) ──────────────────────────────────────
MAROON     = RGBColor(0x8C, 0x1D, 0x40)   # ASU Maroon
GOLD       = RGBColor(0xFF, 0xC6, 0x27)   # ASU Gold
DARK_BG    = RGBColor(0x2B, 0x0A, 0x1C)   # Deep maroon (header)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF4, 0xF4)
MID_GRAY   = RGBColor(0xDD, 0xDD, 0xDD)
TEXT_DARK  = RGBColor(0x1A, 0x1A, 0x1A)
HIGHLIGHT  = RGBColor(0xFF, 0xEC, 0x6E)   # gold highlight for best results
TEAL       = RGBColor(0x00, 0x7A, 0x87)   # accent for pipeline stages
GREEN_ACT  = RGBColor(0x2E, 0x7D, 0x32)   # deep green for gains
BLUE_ACT   = RGBColor(0x1A, 0x56, 0x8C)   # blue accent
SECTION_BG = RGBColor(0xFC, 0xF8, 0xF0)   # warm off-white section bg

# ── Dimensions (EMU) ───────────────────────────────────────────────────────
IN  = 914400   # 1 inch in EMU
W   = int(36 * IN)
H   = int(48 * IN)

MARGIN      = int(0.22 * IN)
HDR_H       = int(4.35 * IN)
BODY_TOP    = HDR_H + int(0.18 * IN)
BODY_H      = H - BODY_TOP - MARGIN
INNER_W     = W - 2 * MARGIN
COL_GAP     = int(0.22 * IN)
N_COLS      = 3
COL_W       = (INNER_W - (N_COLS - 1) * COL_GAP) // N_COLS

# Column x positions
COL_X = [MARGIN + i * (COL_W + COL_GAP) for i in range(N_COLS)]

SEC_HEADER_H = int(0.52 * IN)
ROW_GAP      = int(0.18 * IN)


# ──────────────────────────────────────────────────────────────────────────
# Low-level helpers
# ──────────────────────────────────────────────────────────────────────────

def _rgb(r, g, b):
    return RGBColor(r, g, b)

def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_width
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, x, y, w, h, text, font_size=Pt(11), bold=False,
                color=TEXT_DARK, align=PP_ALIGN.LEFT, italic=False,
                wrap=True, valign=None):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if valign:
        tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_paragraph(tf, text, font_size=Pt(11), bold=False, color=TEXT_DARK,
                  align=PP_ALIGN.LEFT, italic=False, space_before=Pt(0),
                  space_after=Pt(0), bullet=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def section_box(slide, x, y, w, h, title, title_color=WHITE,
                header_fill=MAROON, body_fill=SECTION_BG,
                border_color=None, header_height=None):
    """Draw a titled section box with colored header bar."""
    hh = header_height or SEC_HEADER_H
    # Header bar
    hdr = add_rect(slide, x, y, w, hh, fill=header_fill,
                   line_color=border_color)
    add_textbox(slide, x + int(0.12*IN), y, w - int(0.2*IN), hh,
                title, font_size=Pt(14.5), bold=True, color=title_color,
                align=PP_ALIGN.LEFT)
    # Body
    body = add_rect(slide, x, y + hh, w, h - hh, fill=body_fill,
                    line_color=MID_GRAY, line_width=Pt(0.5))
    return hh   # returns header height so caller knows where body starts


# ──────────────────────────────────────────────────────────────────────────
# Section content builders
# ──────────────────────────────────────────────────────────────────────────

def draw_motivation(slide, x, y, w, h):
    hh = section_box(slide, x, y, w, h, "  Motivation  ·  Why CRAFT?")
    bx = x + int(0.15*IN)
    by = y + hh + int(0.12*IN)
    bw = w - int(0.3*IN)

    txBox = slide.shapes.add_textbox(bx, by, bw, h - hh - int(0.15*IN))
    tf = txBox.text_frame
    tf.word_wrap = True

    lines = [
        ("Open-Domain Table QA", True, Pt(12), MAROON),
        ("Retrieve the right table from 169K–419K tables, then answer\nnatural language queries from its contents.", False, Pt(10.5), TEXT_DARK),
        ("", False, Pt(5), TEXT_DARK),
        ("The Challenge", True, Pt(12), MAROON),
        ("▸  Sparse (BM25/SPLADE): fast but misses semantic meaning", False, Pt(10.5), TEXT_DARK),
        ("▸  Dense (DTR/BIBERT): accurate but requires costly fine-tuning\n    on each new dataset — limiting adaptability", False, Pt(10.5), TEXT_DARK),
        ("▸  Hybrid methods: strong, yet still dataset-specific", False, Pt(10.5), TEXT_DARK),
        ("", False, Pt(5), TEXT_DARK),
        ("The Gap", True, Pt(12), MAROON),
        ("No training-free, plug-and-play retrieval pipeline achieves\nstate-of-the-art recall at scale.", False, Pt(10.5), TEXT_DARK),
        ("", False, Pt(5), TEXT_DARK),
        ("Research Question", True, Pt(12), MAROON),
        ("Can off-the-shelf pretrained models, combined in a carefully\ndesigned cascade, match costly fine-tuned retrievers?", False, Pt(10.5), BLUE_ACT),
    ]
    first = True
    for text, bold, fsize, clr in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = text
        run.font.size = fsize
        run.font.bold = bold
        run.font.color.rgb = clr


def draw_contributions(slide, x, y, w, h):
    hh = section_box(slide, x, y, w, h, "  Key Contributions")
    bx = x + int(0.15*IN)
    by = y + hh + int(0.12*IN)
    bw = w - int(0.3*IN)

    txBox = slide.shapes.add_textbox(bx, by, bw, h - hh - int(0.15*IN))
    tf = txBox.text_frame
    tf.word_wrap = True

    items = [
        ("① Training-Free Cascade",
         "A modular, zero-shot pipeline that requires no dataset-specific\nfine-tuning, enabling rapid deployment to new domains."),
        ("② State-of-the-Art Retrieval",
         "New SOTA Recall@1 on NQ-Tables; competitive zero-shot\nperformance on OTT-QA across all recall thresholds."),
        ("③ Robustness to Paraphrasing",
         "Three-stage cascade nearly fully recovers recall on\nperturbed queries — where DTR degrades by 8–12 points."),
        ("④ Efficiency via Mini-Tables",
         "Reduces token length by >70% and embedding calls\nby 33× vs. dense baselines, enabling scalable inference."),
    ]
    first = True
    for title, body in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(4)
        run = p.add_run()
        run.text = title
        run.font.size = Pt(11.5)
        run.font.bold = True
        run.font.color.rgb = MAROON

        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = body
        run2.font.size = Pt(10.5)
        run2.font.color.rgb = TEXT_DARK
        p2.space_after = Pt(2)


def draw_pipeline(slide, x, y, w, h):
    hh = section_box(slide, x, y, w, h, "  CRAFT Pipeline  ·  Three-Stage Cascade")
    bx = x + int(0.15*IN)
    by = y + hh + int(0.10*IN)
    bw = w - int(0.3*IN)
    inner_h = h - hh - int(0.15*IN)

    # Pre-processing subsection
    pre_h = int(1.35 * IN)
    pre_bg = add_rect(slide, x + int(0.10*IN), by, bw, pre_h,
                      fill=RGBColor(0xE8, 0xF4, 0xFD), line_color=TEAL, line_width=Pt(0.8))
    add_textbox(slide, bx, by + int(0.04*IN), bw, int(0.38*IN),
                "PRE-PROCESSING  (Gemini Flash 1.5)",
                font_size=Pt(11), bold=True, color=TEAL)
    add_textbox(slide, bx, by + int(0.38*IN), bw, int(1.0*IN),
                "▸ Query Decomposition: break complex queries into sub-questions\n"
                "▸ Table Enrichment: generate descriptive titles + summaries\n"
                "▸ Row Ranking: rank rows by semantic relevance (all-mpnet-base-v2)",
                font_size=Pt(10), color=TEXT_DARK)

    # Stage boxes
    stage_colors = [
        (RGBColor(0xFF, 0xF3, 0xE0), RGBColor(0xE6, 0x5C, 0x00)),  # Stage 1 - orange
        (RGBColor(0xE8, 0xF5, 0xE9), RGBColor(0x2E, 0x7D, 0x32)),  # Stage 2 - green
        (RGBColor(0xE8, 0xEA, 0xF6), RGBColor(0x28, 0x35, 0x93)),  # Stage 3 - blue
    ]
    stage_data = [
        ("STAGE 1  ·  Sparse Retrieval",
         "SPLADE lexical-expansion retriever applied over the full table\ncorpus (169K NQ-Tables / 419K OTT-QA tables). Uses enriched\ntable titles + descriptions + query sub-questions.\n→  Outputs N₁ candidate tables (top 5,000)"),
        ("STAGE 2  ·  Dense Semantic Pruning",
         "Sentence Transformer (all-mpnet-base-v2 / Jina v3) re-ranks\nN₁ tables using mini-tables (top 5 rows per table). Sublinear\nANN search over pre-encoded row embeddings.\n→  Outputs N₂ tables (top 100–300)"),
        ("STAGE 3  ·  Neural Re-ranking",
         "High-precision semantic re-ranker (text-embedding-3-small /\ntext-embedding-3-large / gemini-embedding-001) selects the\nfinal top-K mini-tables for answer generation.\n→  Outputs top-K tables (K ∈ {1, 3, 5, 8, 10})"),
    ]

    sy = by + pre_h + int(0.12*IN)
    stage_h = (inner_h - pre_h - int(0.12*IN) - int(0.05*IN)) // 3 - int(0.08*IN)

    for i, ((bg, accent), (stitle, sbody)) in enumerate(zip(stage_colors, stage_data)):
        sx = x + int(0.10*IN)
        add_rect(slide, sx, sy, bw, stage_h, fill=bg,
                 line_color=accent, line_width=Pt(1.0))
        add_textbox(slide, bx, sy + int(0.05*IN), bw, int(0.38*IN),
                    stitle, font_size=Pt(11), bold=True, color=accent)
        add_textbox(slide, bx, sy + int(0.38*IN), bw, stage_h - int(0.42*IN),
                    sbody, font_size=Pt(9.8), color=TEXT_DARK)
        sy += stage_h + int(0.08*IN)


def draw_retrieval_nq(slide, x, y, w, h):
    hh = section_box(slide, x, y, w, h, "  Retrieval Performance  ·  NQ-Tables")
    bx = x + int(0.10*IN)
    by = y + hh + int(0.10*IN)
    bw = w - int(0.20*IN)

    # Intro text
    intro_h = int(0.60*IN)
    add_textbox(slide, bx, by, bw, intro_h,
                "All models except CRAFT are trained on NQ-Tables (169,898 tables). "
                "R = Recall.  Best in bold; second-best underlined.",
                font_size=Pt(9.5), color=TEXT_DARK, italic=True)

    # Table
    ty = by + intro_h
    table_h = h - hh - intro_h - int(0.15*IN)

    rows = [
        # (label, R@1, R@10, R@50, category, bold, highlight)
        ("Sparse Retrieval", "", "", "", "cat", False, False),
        ("BM25",            "18.49", "36.94", "52.61", "sparse", False, False),
        ("SPLADE",          "39.84", "83.33", "94.65", "sparse", False, False),
        ("Dense Retrieval", "", "", "", "cat", False, False),
        ("BIBERT",          "43.78", "82.25", "93.71", "dense", False, False),
        ("DPR",             "45.32", "85.84", "95.44", "dense", False, False),
        ("DTR",             "32.62", "75.86", "89.77", "dense", False, False),
        ("SSDR",            "45.47", "84.00", "95.05", "dense", False, False),
        ("Hybrid Retrieval","", "", "", "cat", False, False),
        ("BIBERT+SPLADE",   "43.67", "86.72", "95.62", "hybrid", False, False),
        ("THYME",           "48.55", "86.38", "96.08", "hybrid", False, False),
        ("CRAFT Stage-wise","", "", "", "cat", False, False),
        ("Stage 1",         "34.38", "72.90", "91.62", "craft_s", False, False),
        ("+ Stage 2",       "36.65", "82.91", "96.08", "craft_s", False, False),
        ("+ Stage 3",       "41.13", "87.16", "96.84", "craft_s", False, False),
        ("CRAFT (full)",    "49.84", "86.83", "97.17", "craft", True, True),
    ]

    n_rows = len(rows)
    row_h = min(int(table_h / (n_rows + 1)), int(0.30 * IN))
    col_ws = [int(bw * 0.42), int(bw * 0.19), int(bw * 0.19), int(bw * 0.20)]
    headers = ["Model", "R@1", "R@10", "R@50"]
    cat_fills = {
        "cat":     RGBColor(0xCC, 0xCC, 0xCC),
        "sparse":  RGBColor(0xF5, 0xF5, 0xF5),
        "dense":   RGBColor(0xF0, 0xF4, 0xFF),
        "hybrid":  RGBColor(0xF0, 0xFF, 0xF4),
        "craft_s": RGBColor(0xFF, 0xF8, 0xE8),
        "craft":   HIGHLIGHT,
    }

    # Header row
    cx = bx
    for ci, (hdr, cw) in enumerate(zip(headers, col_ws)):
        add_rect(slide, cx, ty, cw, row_h, fill=MAROON)
        add_textbox(slide, cx + int(0.03*IN), ty, cw, row_h,
                    hdr, font_size=Pt(10), bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
        cx += cw

    ty += row_h
    for label, r1, r10, r50, cat, bold, hl in rows:
        fill = cat_fills.get(cat, LIGHT_GRAY)
        cx = bx
        vals = [label, r1, r10, r50]
        for ci, (val, cw) in enumerate(zip(vals, col_ws)):
            add_rect(slide, cx, ty, cw, row_h, fill=fill,
                     line_color=MID_GRAY, line_width=Pt(0.3))
            txt_color = MAROON if (cat == "cat") else \
                        RGBColor(0x8C, 0x1D, 0x40) if (cat == "craft" and bold) else TEXT_DARK
            add_textbox(slide, cx + int(0.03*IN), ty, cw - int(0.03*IN), row_h,
                        val, font_size=Pt(9.5 if cat != "cat" else 9),
                        bold=(bold or cat == "cat"),
                        color=txt_color,
                        align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
            cx += cw
        ty += row_h

    # SOTA callout
    callout_y = ty + int(0.08*IN)
    callout_h = int(0.52*IN)
    remaining = (y + h) - callout_y - int(0.05*IN)
    if remaining > int(0.40*IN):
        add_rect(slide, bx, callout_y, bw, min(callout_h, remaining),
                 fill=HIGHLIGHT, line_color=MAROON, line_width=Pt(1.0))
        add_textbox(slide, bx + int(0.08*IN), callout_y, bw - int(0.1*IN),
                    min(callout_h, remaining),
                    "CRAFT achieves new SOTA  R@1 = 49.84  on NQ-Tables "
                    "(trained models excluded from comparison)",
                    font_size=Pt(10), bold=True, color=MAROON)


def draw_retrieval_ottqa(slide, x, y, w, h):
    hh = section_box(slide, x, y, w, h, "  Retrieval Performance  ·  OTT-QA (Zero-Shot)")
    bx = x + int(0.10*IN)
    by = y + hh + int(0.10*IN)
    bw = w - int(0.20*IN)

    intro_h = int(0.60*IN)
    add_textbox(slide, bx, by, bw, intro_h,
                "OTT-QA (419,183 tables) requires multi-hop reasoning over tables + passages. "
                "CRAFT is zero-shot — all other methods are fine-tuned.",
                font_size=Pt(9.5), color=TEXT_DARK, italic=True)

    ty = by + intro_h
    table_h = h - hh - intro_h - int(0.65*IN)

    rows = [
        ("Sparse Retrieval", "", "", "", "cat"),
        ("BM25",             "23.98", "51.94", "69.11", "sparse"),
        ("SPLADE",           "62.74", "89.52", "95.21", "sparse"),
        ("Dense Retrieval",  "", "", "", "cat"),
        ("BIBERT",           "56.82", "86.50", "94.26", "dense"),
        ("DPR",              "53.43", "85.95", "93.22", "dense"),
        ("SSDR",             "56.96", "86.22", "93.55", "dense"),
        ("Hybrid Retrieval", "", "", "", "cat"),
        ("BIBERT+SPLADE",    "59.49", "86.81", "94.67", "hybrid"),
        ("BIBERT+SPLADE",    "64.72", "91.01", "96.34", "hybrid"),  # underline
        ("THYME",            "66.67", "91.10", "96.16", "hybrid"),
        ("CRAFT Stage-wise", "", "", "", "cat"),
        ("Stage 1",          "40.74", "60.89", "88.17", "craft_s"),
        ("+ Stage 2",        "47.33", "87.35", "91.82", "craft_s"),
        ("CRAFT (full)",     "55.56", "89.88", "96.07", "craft"),
    ]

    n_rows = len(rows)
    row_h = min(int(table_h / (n_rows + 1)), int(0.30 * IN))
    col_ws = [int(bw * 0.42), int(bw * 0.19), int(bw * 0.19), int(bw * 0.20)]
    headers = ["Model", "R@1", "R@10", "R@50"]
    cat_fills = {
        "cat":     RGBColor(0xCC, 0xCC, 0xCC),
        "sparse":  RGBColor(0xF5, 0xF5, 0xF5),
        "dense":   RGBColor(0xF0, 0xF4, 0xFF),
        "hybrid":  RGBColor(0xF0, 0xFF, 0xF4),
        "craft_s": RGBColor(0xFF, 0xF8, 0xE8),
        "craft":   HIGHLIGHT,
    }

    cx_start = bx
    for ci, (hdr, cw) in enumerate(zip(headers, col_ws)):
        add_rect(slide, cx_start, ty, cw, row_h, fill=MAROON)
        add_textbox(slide, cx_start + int(0.03*IN), ty, cw, row_h,
                    hdr, font_size=Pt(10), bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
        cx_start += cw

    ty += row_h
    for label, r1, r10, r50, cat in rows:
        fill = cat_fills.get(cat, LIGHT_GRAY)
        cx = bx
        vals = [label, r1, r10, r50]
        bold = (cat == "craft" or cat == "cat")
        txt_color = MAROON if cat in ("cat", "craft") else TEXT_DARK
        for ci, (val, cw) in enumerate(zip(vals, col_ws)):
            add_rect(slide, cx, ty, cw, row_h, fill=fill,
                     line_color=MID_GRAY, line_width=Pt(0.3))
            add_textbox(slide, cx + int(0.03*IN), ty, cw - int(0.03*IN), row_h,
                        val, font_size=Pt(9.5 if cat != "cat" else 9),
                        bold=bold, color=txt_color,
                        align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
            cx += cw
        ty += row_h

    # Callout
    callout_y = ty + int(0.08*IN)
    remaining = (y + h) - callout_y - int(0.05*IN)
    if remaining > int(0.35*IN):
        add_rect(slide, bx, callout_y, bw, min(int(0.52*IN), remaining),
                 fill=RGBColor(0xE8, 0xF4, 0xFD), line_color=TEAL, line_width=Pt(1.0))
        add_textbox(slide, bx + int(0.08*IN), callout_y,
                    bw - int(0.1*IN), min(int(0.52*IN), remaining),
                    "Zero-shot CRAFT approaches fine-tuned THYME at R@10 & R@50 "
                    "— with no dataset-specific training",
                    font_size=Pt(10), bold=True, color=TEAL)


def draw_e2e_results(slide, x, y, w, h):
    hh = section_box(slide, x, y, w, h, "  End-to-End QA Performance  ·  NQ-Tables (F1)")
    bx = x + int(0.10*IN)
    by = y + hh + int(0.10*IN)
    bw = w - int(0.20*IN)

    intro_h = int(0.52*IN)
    add_textbox(slide, bx, by, bw, intro_h,
                "F1 scores with varying number of retrieved tables (n). "
                "Models: Mistral-7B-Instruct, Llama3-8B-Instruct, Qwen2.5-7B-Instruct.",
                font_size=Pt(9.5), color=TEXT_DARK, italic=True)

    ty = by + intro_h
    table_h = h - hh - intro_h - int(0.10*IN)
    row_h = min(int(table_h / 8), int(0.31 * IN))

    # Header
    col_labels = ["Retriever", "n=1\nMistral", "Llama3", "Qwen",
                  "n=3\nMistral", "Llama3", "Qwen",
                  "n=5\nMistral", "Llama3", "Qwen"]
    col_ws_raw = [0.26, 0.082, 0.082, 0.082,
                  0.082, 0.082, 0.082,
                  0.082, 0.082, 0.082]
    col_ws = [int(bw * p) for p in col_ws_raw]

    # Group header row 1 — n=1, n=3, n=5
    grp_labels = [("", 1), ("n = 1", 3), ("n = 3", 3), ("n = 5", 3)]
    cx = bx
    for glabel, span in grp_labels:
        gw = sum(col_ws[:span]) if glabel == "" else sum(col_ws[1:1+span]) if glabel == "n = 1" \
             else sum(col_ws[4:4+span]) if glabel == "n = 3" else sum(col_ws[7:7+span])
        if glabel == "":
            gw = col_ws[0]
        fill_c = MAROON if glabel else RGBColor(0x6B, 0x15, 0x30)
        add_rect(slide, cx, ty, gw, row_h, fill=fill_c)
        add_textbox(slide, cx, ty, gw, row_h,
                    glabel, font_size=Pt(9.5), bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        cx += gw

    ty += row_h
    # Sub-header
    cx = bx
    sub_headers = ["Retriever", "Mis.", "Ll3", "Qwn", "Mis.", "Ll3", "Qwn", "Mis.", "Ll3", "Qwn"]
    for sh, cw in zip(sub_headers, col_ws):
        add_rect(slide, cx, ty, cw, row_h, fill=RGBColor(0xAA, 0x28, 0x50))
        add_textbox(slide, cx, ty, cw, row_h,
                    sh, font_size=Pt(8.5), bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        cx += cw
    ty += row_h

    data_rows = [
        ("BIBERT",       "32.93","32.66","34.80","34.6","33.16","37.92","35.30","34.28","37.01", False),
        ("SPLADE",       "29.61","32.07","31.90","35.42","37.17","35.79","34.95","33.88","37.35", False),
        ("BIBERT+SPLADE","32.67","32.66","33.24","33.59","33.66","36.76","35.71","33.92","37.02", False),
        ("THYME",        "35.48","36.14","37.28","37.59","39.16","40.28","37.20","39.29","41.20", False),
        ("CRAFT",        "39.13","38.31","39.73","45.28","40.76","43.52","44.53","40.55","46.49", True),
    ]

    row_fills = [RGBColor(0xF5,0xF5,0xF5), RGBColor(0xF0,0xF4,0xFF),
                 RGBColor(0xF0,0xFF,0xF4), RGBColor(0xFF,0xF8,0xE8), HIGHLIGHT]

    for ri, (label, *vals_b) in enumerate(data_rows):
        bold = vals_b[-1]
        vals = vals_b[:-1]
        fill = row_fills[ri]
        cx = bx
        for ci, (val, cw) in enumerate(zip([label]+vals, col_ws)):
            add_rect(slide, cx, ty, cw, row_h, fill=fill,
                     line_color=MID_GRAY, line_width=Pt(0.3))
            txt_c = MAROON if bold else TEXT_DARK
            add_textbox(slide, cx + int(0.02*IN), ty, cw - int(0.02*IN), row_h,
                        val, font_size=Pt(9 if ci == 0 else 8.5),
                        bold=bold, color=txt_c,
                        align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
            cx += cw
        ty += row_h


def draw_robustness(slide, x, y, w, h):
    hh = section_box(slide, x, y, w, h, "  Robustness to Query Paraphrasing",
                     header_fill=RGBColor(0x1A, 0x56, 0x8C))
    bx = x + int(0.10*IN)
    by = y + hh + int(0.10*IN)
    bw = w - int(0.20*IN)

    add_textbox(slide, bx, by, bw, int(0.55*IN),
                "Perturbed queries generated by Gemini 2.5 Flash. "
                "Δ = change vs. original query (negative = degradation).",
                font_size=Pt(9.5), color=TEXT_DARK, italic=True)

    ty = by + int(0.55*IN)
    row_h = int(0.29*IN)
    col_ws = [int(bw * 0.32), int(bw * 0.17), int(bw * 0.17), int(bw * 0.17), int(bw * 0.17)]
    headers = ["Model", "R@1", "R@10", "R@50", "Δ Avg."]

    cx = bx
    for ci, (hdr, cw) in enumerate(zip(headers, col_ws)):
        add_rect(slide, cx, ty, cw, row_h, fill=RGBColor(0x1A, 0x56, 0x8C))
        add_textbox(slide, cx + int(0.02*IN), ty, cw, row_h,
                    hdr, font_size=Pt(9.5), bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
        cx += cw
    ty += row_h

    perturbed_rows = [
        ("Original Query", "", "", "", "", "cat"),
        ("DTR (M)",   "38.74","75.73","88.36","–", "orig"),
        ("CRAFT",     "41.13","87.16","96.84","–", "craft"),
        ("Perturbed Query","","","","","cat"),
        ("DTR (M)",   "27.09","66.27","84.33","–8.38","bad"),
        ("DTR (M)+HN","37.87","76.06","84.43","–5.80","bad"),
        ("DTR (S)+HN","35.04","70.73","85.85","–6.64","bad"),
        ("CRAFT",     "41.02","86.83","96.08","–0.04","craft"),
    ]
    fills = {
        "cat":   RGBColor(0xCC, 0xCC, 0xCC),
        "orig":  RGBColor(0xF0, 0xF4, 0xFF),
        "bad":   RGBColor(0xFF, 0xF0, 0xF0),
        "craft": HIGHLIGHT,
    }
    for label, r1, r10, r50, delta, cat in perturbed_rows:
        fill = fills[cat]
        cx = bx
        bold = (cat in ("cat","craft"))
        for ci, (val, cw) in enumerate(zip([label,r1,r10,r50,delta], col_ws)):
            add_rect(slide, cx, ty, cw, row_h, fill=fill,
                     line_color=MID_GRAY, line_width=Pt(0.3))
            txt_c = MAROON if cat == "craft" else (RGBColor(0xAA,0x00,0x00) if (cat=="bad" and ci==4 and val) else TEXT_DARK)
            add_textbox(slide, cx + int(0.02*IN), ty, cw - int(0.02*IN), row_h,
                        val, font_size=Pt(9 if ci==0 else 8.5),
                        bold=bold, color=txt_c,
                        align=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER)
            cx += cw
        ty += row_h

    # Callout
    cally = ty + int(0.08*IN)
    rem = (y + h) - cally - int(0.05*IN)
    if rem > int(0.35*IN):
        add_rect(slide, bx, cally, bw, min(int(0.52*IN), rem),
                 fill=HIGHLIGHT, line_color=MAROON, line_width=Pt(0.8))
        add_textbox(slide, bx + int(0.08*IN), cally,
                    bw - int(0.1*IN), min(int(0.52*IN), rem),
                    "CRAFT degrades by only –0.04 points on perturbed queries "
                    "vs. –5.8 to –11.8 for DTR variants",
                    font_size=Pt(10), bold=True, color=MAROON)


def draw_efficiency(slide, x, y, w, h):
    hh = section_box(slide, x, y, w, h, "  Efficiency  ·  Mini-Tables & Embedding Calls",
                     header_fill=GREEN_ACT)
    bx = x + int(0.10*IN)
    by = y + hh + int(0.10*IN)
    bw = w - int(0.20*IN)

    txBox = slide.shapes.add_textbox(bx, by, bw, h - hh - int(0.15*IN))
    tf = txBox.text_frame
    tf.word_wrap = True

    lines = [
        ("Token Reduction via Mini-Tables", True, Pt(11), MAROON),
        ("Mini-tables (top 5 rows) reduce context length by >70% vs. full tables:", False, Pt(10), TEXT_DARK),
        ("  • Llama3: 1782 → 363 tokens (k=1)   •  Mistral: 3314 → 477", False, Pt(10), TEXT_DARK),
        ("  • Qwen2.5: 1951 → 394 tokens (k=1)   reduces noise for downstream LLM", False, Pt(10), TEXT_DARK),
        ("", False, Pt(4), TEXT_DARK),
        ("Embedding Call Reduction", True, Pt(11), MAROON),
        ("CRAFT executes only 3 query embeddings + 5,100 mini-table\n"
         "embeddings per query at inference time:", False, Pt(10), TEXT_DARK),
        ("  Stage 1: 1 query encoding + inverted-index lookup", False, Pt(10), TEXT_DARK),
        ("  Stage 2: 1 ANN search over pre-encoded row embeddings", False, Pt(10), TEXT_DARK),
        ("  Stage 3: 1 query embed + top-100 mini-table re-ranking", False, Pt(10), TEXT_DARK),
        ("", False, Pt(3), TEXT_DARK),
        ("Vs. fully dense baseline: 169,898 table embedding calls", False, Pt(10),
         RGBColor(0xAA,0x00,0x00)),
        ("CRAFT reduces embedding calls by 33× without accuracy loss.", True, Pt(10.5), GREEN_ACT),
        ("", False, Pt(4), TEXT_DARK),
        ("Adaptability", True, Pt(11), MAROON),
        ("Stage 3 re-ranker can be swapped for any new embedding model\n"
         "(e.g., text-embedding-3-large, gemini-embedding-001) with zero\n"
         "architectural changes — future-proof to new embedding releases.", False, Pt(10), TEXT_DARK),
    ]

    first = True
    for text, bold, fsize, clr in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = text
        run.font.size = fsize
        run.font.bold = bold
        run.font.color.rgb = clr


def draw_conclusion(slide, x, y, w, h):
    hh = section_box(slide, x, y, w, h, "  Conclusion",
                     header_fill=DARK_BG)
    bx = x + int(0.10*IN)
    by = y + hh + int(0.08*IN)
    bw = w - int(0.20*IN)

    txBox = slide.shapes.add_textbox(bx, by, bw, h - hh - int(0.12*IN))
    tf = txBox.text_frame
    tf.word_wrap = True

    lines = [
        ("CRAFT is a modular, training-free cascaded retrieval framework for "
         "Open-Domain Table QA that achieves state-of-the-art retrieval without "
         "any dataset-specific fine-tuning.", False, Pt(10.5), TEXT_DARK),
        ("", False, Pt(3), TEXT_DARK),
        ("Key Takeaways", True, Pt(11), MAROON),
        ("▸  New SOTA Recall@1 on NQ-Tables (49.84)", False, Pt(10.5), TEXT_DARK),
        ("▸  Competitive zero-shot performance on OTT-QA", False, Pt(10.5), TEXT_DARK),
        ("▸  33× fewer embedding calls vs. dense baselines", False, Pt(10.5), TEXT_DARK),
        ("▸  Robust to query paraphrasing (–0.04 Δ)", False, Pt(10.5), TEXT_DARK),
        ("▸  Plug-and-play: swap any embedding model at Stage 3", False, Pt(10.5), TEXT_DARK),
        ("", False, Pt(3), TEXT_DARK),
        ("Impact: Bridges fine-tuned & lightweight retrieval — scalable, "
         "interpretable, and adaptable to evolving embedding ecosystems.", False, Pt(10.5), BLUE_ACT),
    ]

    first = True
    for text, bold, fsize, clr in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = text
        run.font.size = fsize
        run.font.bold = bold
        run.font.color.rgb = clr


def draw_ablation_stagewise(slide, x, y, w, h):
    """Stage-wise ablation bar chart (text-based) showing progressive gains."""
    hh = section_box(slide, x, y, w, h, "  Ablation  ·  Stage-wise Gains",
                     header_fill=RGBColor(0x5C, 0x1A, 0x5C))
    bx = x + int(0.10*IN)
    by = y + hh + int(0.10*IN)
    bw = w - int(0.20*IN)

    add_textbox(slide, bx, by, bw, int(0.40*IN),
                "Recall@10 improvements at each stage (NQ-Tables):",
                font_size=Pt(10), bold=True, color=TEXT_DARK)

    # Bar visualization
    bars = [
        ("Stage 1\n(SPLADE)",      72.90, RGBColor(0xE6,0x5C,0x00)),
        ("+ Stage 2\n(Dense)",     82.91, RGBColor(0x2E,0x7D,0x32)),
        ("+ Stage 3\n(Rerank)",    87.16, RGBColor(0x28,0x35,0x93)),
        ("CRAFT\n(full)",          86.83, MAROON),
    ]

    bar_area_y = by + int(0.42*IN)
    bar_area_h = h - hh - int(0.55*IN)
    max_val = 100.0
    bar_w = int((bw - int(0.3*IN)) / len(bars)) - int(0.08*IN)
    bbar_x = bx + int(0.15*IN)

    for i, (label, val, color) in enumerate(bars):
        bx_i = bbar_x + i * (bar_w + int(0.08*IN))
        bar_fill_h = int((val / max_val) * (bar_area_h - int(0.50*IN)))
        bar_y = bar_area_y + (bar_area_h - int(0.50*IN)) - bar_fill_h

        # Background bar (gray)
        add_rect(slide, bx_i, bar_area_y, bar_w,
                 bar_area_h - int(0.50*IN),
                 fill=RGBColor(0xEE,0xEE,0xEE))
        # Colored fill
        add_rect(slide, bx_i, bar_y, bar_w, bar_fill_h,
                 fill=color)
        # Value label
        add_textbox(slide, bx_i - int(0.05*IN), bar_y - int(0.32*IN),
                    bar_w + int(0.10*IN), int(0.30*IN),
                    f"{val:.1f}",
                    font_size=Pt(9.5), bold=True, color=color,
                    align=PP_ALIGN.CENTER)
        # X label
        label_y = bar_area_y + bar_area_h - int(0.50*IN) + int(0.04*IN)
        add_textbox(slide, bx_i - int(0.03*IN), label_y,
                    bar_w + int(0.06*IN), int(0.45*IN),
                    label, font_size=Pt(8.5), color=TEXT_DARK,
                    align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────────
# Main builder
# ──────────────────────────────────────────────────────────────────────────

def build_poster():
    prs = Presentation()
    prs.slide_width  = Emu(W)
    prs.slide_height = Emu(H)

    blank_layout = prs.slide_layouts[6]   # completely blank
    slide = prs.slides.add_slide(blank_layout)

    # ── Full background ────────────────────────────────────────────────────
    add_rect(slide, 0, 0, W, H, fill=RGBColor(0xF0, 0xEE, 0xEB))

    # ── HEADER ────────────────────────────────────────────────────────────
    add_rect(slide, 0, 0, W, HDR_H, fill=DARK_BG)
    # Subtle gold accent stripe
    add_rect(slide, 0, HDR_H - int(0.07*IN), W, int(0.07*IN), fill=GOLD)

    # Title
    title_tb = slide.shapes.add_textbox(
        int(0.30*IN), int(0.18*IN),
        int(25.5*IN), int(1.85*IN))
    tf = title_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "CRAFT: Training-Free Cascaded Retrieval for Tabular QA"
    run.font.size = Pt(46)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Authors
    authors_tb = slide.shapes.add_textbox(
        int(0.30*IN), int(2.10*IN),
        int(25.5*IN), int(0.70*IN))
    tf2 = authors_tb.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = ("*Adarsh Singh¹   *Kushal Raj Bhandari²   Jianxi Gao³   "
                 "Soham Dan²   †Vivek Gupta¹")
    run2.font.size = Pt(20)
    run2.font.color.rgb = GOLD

    # Affiliations
    aff_tb = slide.shapes.add_textbox(
        int(0.30*IN), int(2.85*IN),
        int(25.5*IN), int(0.55*IN))
    tf3 = aff_tb.text_frame
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = ("¹Arizona State University    ²Rensselaer Polytechnic Institute    "
                 "³Microsoft    *Equal contribution    †Primary supervisor")
    run3.font.size = Pt(16)
    run3.font.color.rgb = RGBColor(0xDD, 0xCC, 0xBB)

    # Venue badge
    badge_x = int(26.2*IN)
    badge_w = int(5.5*IN)
    add_rect(slide, badge_x, int(0.20*IN), badge_w, int(1.10*IN),
             fill=MAROON, line_color=GOLD, line_width=Pt(2))
    add_textbox(slide, badge_x + int(0.10*IN), int(0.25*IN),
                badge_w - int(0.15*IN), int(1.05*IN),
                "ACL 2026  ·  Vienna, Austria",
                font_size=Pt(18), bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)

    add_textbox(slide, badge_x + int(0.10*IN), int(1.40*IN),
                badge_w - int(0.15*IN), int(0.55*IN),
                "arXiv: 2505.14984",
                font_size=Pt(15), bold=False, color=GOLD,
                align=PP_ALIGN.CENTER)

    add_textbox(slide, badge_x + int(0.10*IN), int(2.05*IN),
                badge_w - int(0.15*IN), int(0.55*IN),
                "asing725@asu.edu",
                font_size=Pt(13), color=RGBColor(0xDD,0xCC,0xBB),
                align=PP_ALIGN.CENTER)

    add_textbox(slide, badge_x + int(0.10*IN), int(2.65*IN),
                badge_w - int(0.15*IN), int(0.55*IN),
                "bhandk@rpi.edu",
                font_size=Pt(13), color=RGBColor(0xDD,0xCC,0xBB),
                align=PP_ALIGN.CENTER)

    # ── LOGOS ──────────────────────────────────────────────────────────────
    try:
        pic = slide.shapes.add_picture(
            "/mnt/data1/asing725/ACL/CRAFT/static/images/asu_logo.png",
            int(31.9*IN), int(3.30*IN), int(1.80*IN), int(0.70*IN))
    except Exception:
        pass
    try:
        pic2 = slide.shapes.add_picture(
            "/mnt/data1/asing725/ACL/CRAFT/static/images/coral_logo.jpeg",
            int(33.85*IN), int(3.20*IN), int(1.90*IN), int(0.90*IN))
    except Exception:
        pass

    # ── BODY LAYOUT ────────────────────────────────────────────────────────
    # Row heights
    ROW1_H = int(11.2 * IN)
    ROW2_H = int(11.8 * IN)
    ROW3_H = int( 9.8 * IN)
    ROW4_H = H - BODY_TOP - ROW1_H - ROW2_H - ROW3_H - 3*ROW_GAP - MARGIN

    r1y = BODY_TOP
    r2y = r1y + ROW1_H + ROW_GAP
    r3y = r2y + ROW2_H + ROW_GAP
    r4y = r3y + ROW3_H + ROW_GAP

    # ── ROW 1 ──────────────────────────────────────────────────────────────
    # Col 0: Motivation
    draw_motivation(slide, COL_X[0], r1y, COL_W, ROW1_H)
    # Col 1: CRAFT overview figure (large, center)
    ov_hh = section_box(slide, COL_X[1], r1y, COL_W, ROW1_H,
                        "  CRAFT Framework Overview")
    try:
        img_y = r1y + ov_hh + int(0.10*IN)
        img_h = ROW1_H - ov_hh - int(0.15*IN)
        img_w = COL_W - int(0.20*IN)
        slide.shapes.add_picture(
            "/mnt/data1/asing725/ACL/CRAFT/static/images/craft_overview.png",
            COL_X[1] + int(0.10*IN), img_y, img_w, img_h)
    except Exception as e:
        add_textbox(slide, COL_X[1]+int(0.15*IN), r1y+ov_hh+int(0.2*IN),
                    COL_W-int(0.3*IN), int(1.0*IN),
                    f"[craft_overview.png — {e}]", font_size=Pt(9), color=TEXT_DARK)
    # Col 2: Key contributions
    draw_contributions(slide, COL_X[2], r1y, COL_W, ROW1_H)

    # ── ROW 2 ──────────────────────────────────────────────────────────────
    # Col 0: Three-stage pipeline
    draw_pipeline(slide, COL_X[0], r2y, COL_W, ROW2_H)
    # Col 1: NQ-Tables retrieval
    draw_retrieval_nq(slide, COL_X[1], r2y, COL_W, ROW2_H)
    # Col 2: OTT-QA retrieval
    draw_retrieval_ottqa(slide, COL_X[2], r2y, COL_W, ROW2_H)

    # ── ROW 3 ──────────────────────────────────────────────────────────────
    # Col 0+1 span (2/3 width): End-to-end results
    e2e_w = COL_W * 2 + COL_GAP
    draw_e2e_results(slide, COL_X[0], r3y, e2e_w, ROW3_H)
    # Col 2: Ablation / stage-wise
    draw_ablation_stagewise(slide, COL_X[2], r3y, COL_W, ROW3_H)

    # ── ROW 4 ──────────────────────────────────────────────────────────────
    # Col 0: Robustness
    draw_robustness(slide, COL_X[0], r4y, COL_W, ROW4_H)
    # Col 1: Efficiency
    draw_efficiency(slide, COL_X[1], r4y, COL_W, ROW4_H)
    # Col 2: Conclusion
    draw_conclusion(slide, COL_X[2], r4y, COL_W, ROW4_H)

    # ── FOOTER ─────────────────────────────────────────────────────────────
    footer_y = H - int(0.28*IN)
    add_rect(slide, 0, footer_y, W, int(0.28*IN), fill=DARK_BG)
    add_textbox(slide, int(0.3*IN), footer_y, int(20*IN), int(0.26*IN),
                "CRAFT  ·  Training-Free Cascaded Retrieval for Tabular QA  ·  ACL 2026",
                font_size=Pt(11), color=RGBColor(0xCC,0xBB,0xAA),
                align=PP_ALIGN.LEFT)
    add_textbox(slide, int(16*IN), footer_y, int(19.5*IN), int(0.26*IN),
                "arXiv: 2505.14984  ·  asing725@asu.edu  ·  bhandk@rpi.edu",
                font_size=Pt(11), color=GOLD,
                align=PP_ALIGN.RIGHT)

    out = "/mnt/data1/asing725/ACL/CRAFT/CRAFT_poster.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slide size: {prs.slide_width.inches:.1f}\" x {prs.slide_height.inches:.1f}\"")


if __name__ == "__main__":
    build_poster()
