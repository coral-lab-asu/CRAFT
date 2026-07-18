#!/usr/bin/env python3
"""
generate_poster_pptx.py — CRAFT ACL 2026 poster as a single PPTX slide.

Slide size: 36" × 48" portrait  (standard conference poster)
Run:  python3 generate_poster_pptx.py
Output: CRAFT_poster.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Paths ────────────────────────────────────────────────────────────────────
BASE   = os.path.dirname(__file__)
IMG    = os.path.join(BASE, 'static', 'images')
OUT    = os.path.join(BASE, 'CRAFT_poster.pptx')

# ── Colours ──────────────────────────────────────────────────────────────────
MAROON   = RGBColor(0x8C, 0x1D, 0x40)
GOLD     = RGBColor(0xFF, 0xC6, 0x27)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DARK     = RGBColor(0x1A, 0x1A, 0x1A)
LGRAY    = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
BLUE     = RGBColor(0x0D, 0x7C, 0xC5)
GREEN    = RGBColor(0x2E, 0x7D, 0x32)
ORANGE   = RGBColor(0xFF, 0x6B, 0x35)
CREAM    = RGBColor(0xFF, 0xF8, 0xE1)

# ── Slide setup (36" × 48") ──────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(36)
prs.slide_height = Inches(48)

slide_layout = prs.slide_layouts[6]   # blank
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes

# ── Helper functions ──────────────────────────────────────────────────────────

def add_rect(l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    """Add a filled rectangle. Returns the shape."""
    shape = shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line
    else:
        shape.line.fill.background()
    return shape

def add_label(text, l, t, w, h,
              size=Pt(22), bold=False, color=WHITE,
              align=PP_ALIGN.LEFT, italic=False,
              wrap=True, valign='top'):
    """Add a text box with a single paragraph run."""
    txBox = shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if valign == 'middle':
        tf.vertical_anchor = 3   # MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def section_header(title, l, t, w, h=Inches(0.45)):
    """Draw a maroon section header bar with white title text."""
    add_rect(l, t, w, h, fill=MAROON)
    add_label(title, l + Inches(0.12), t + Inches(0.04),
              w - Inches(0.2), h - Inches(0.06),
              size=Pt(28), bold=True, color=WHITE,
              align=PP_ALIGN.LEFT, valign='middle')
    return t + h

def add_image(filename, l, t, w, h):
    img_path = os.path.join(IMG, filename)
    if os.path.exists(img_path):
        shapes.add_picture(img_path, l, t, w, h)
    else:
        box = add_rect(l, t, w, h, fill=LGRAY)
        add_label(f'[{filename}]', l, t, w, h,
                  size=Pt(18), color=DARK, align=PP_ALIGN.CENTER)

def bullet_box(lines, l, t, w, h, size=Pt(21), gap=Pt(4)):
    """Multi-line bullet text box. `lines` = list of (bullet_char, text, bold)."""
    txBox = shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (bul, txt, bld) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = gap
        run = p.add_run()
        run.text = f'{bul} {txt}'
        run.font.size  = size
        run.font.bold  = bld
        run.font.color.rgb = DARK
    return txBox

def table_shape(data, col_widths, row_height,
                l, t, header_fill=MAROON, alt_fill=LGRAY,
                best_row=None, best_fill=CREAM, best_border=GOLD,
                hdr_size=Pt(20), data_size=Pt(20)):
    """
    Draw a table as grouped rectangles + text.
    data[0] = header row  (list of strings)
    best_row = 0-based data row index (not counting header) to highlight
    """
    n_cols = len(data[0])
    total_w = sum(col_widths)
    for r, row in enumerate(data):
        x = l
        for c, val in enumerate(row):
            cw = col_widths[c]
            if r == 0:
                bg = header_fill
                fc = WHITE
                bld = True
                sz = hdr_size
            elif best_row is not None and r - 1 == best_row:
                bg = best_fill
                fc = DARK
                bld = True
                sz = data_size
            elif r % 2 == 0:
                bg = LGRAY
                fc = DARK
                bld = False
                sz = data_size
            else:
                bg = WHITE
                fc = DARK
                bld = False
                sz = data_size
            add_rect(x, t + r * row_height, cw, row_height,
                     fill=bg,
                     line=best_border if (best_row is not None and r - 1 == best_row) else MID_GRAY,
                     line_w=Pt(2) if (best_row is not None and r - 1 == best_row) else Pt(0.5))
            pad_l = Inches(0.1)
            align = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            add_label(str(val),
                      x + pad_l, t + r * row_height + Inches(0.04),
                      cw - pad_l, row_height - Inches(0.06),
                      size=sz, bold=bld, color=fc if r == 0 else DARK,
                      align=align, valign='middle')
            x += cw

def callout_box(text1, text2, l, t, w, h):
    """Gold callout box with two text lines."""
    add_rect(l, t, w, h, fill=GOLD, line=MAROON, line_w=Pt(3))
    add_label(text1, l + Inches(0.1), t + Inches(0.08),
              w - Inches(0.2), h * 0.52,
              size=Pt(32), bold=True, color=MAROON,
              align=PP_ALIGN.CENTER, valign='middle')
    add_label(text2, l + Inches(0.1), t + h * 0.55,
              w - Inches(0.2), h * 0.42,
              size=Pt(20), bold=False, color=DARK,
              align=PP_ALIGN.CENTER, valign='middle')

def stage_badge(text, l, t, w, h, fill):
    add_rect(l, t, w, h, fill=fill)
    add_label(text, l + Inches(0.05), t + Inches(0.04),
              w - Inches(0.1), h - Inches(0.08),
              size=Pt(20), bold=True, color=WHITE,
              align=PP_ALIGN.CENTER, valign='middle')

# ═══════════════════════════════════════════════════════════════════════════════
# LAYOUT CONSTANTS
# ═══════════════════════════════════════════════════════════════════════════════

M      = Inches(0.25)    # outer margin
GUTTER = Inches(0.22)    # column gutter
C1L    = M
CW     = Inches(11.2)    # column width  (3 cols × 11.2 + 2 × 0.22 + 2×0.25 = 36)
C2L    = C1L + CW + GUTTER
C3L    = C2L + CW + GUTTER
BODY_T = Inches(5.1)     # body starts after header
BODY_B = Inches(45.8)    # body ends before footer
SEC_H  = Inches(0.45)    # section header height
INNER  = Inches(0.12)    # inner horizontal pad

# ═══════════════════════════════════════════════════════════════════════════════
# HEADER  (0 → 5.1")
# ═══════════════════════════════════════════════════════════════════════════════

add_rect(0, 0, prs.slide_width, Inches(5.1), fill=MAROON)

# Title
add_label(
    'CRAFT: Training-Free Cascaded Retrieval for Tabular QA',
    Inches(0.3), Inches(0.25),
    Inches(29.5), Inches(2.2),
    size=Pt(68), bold=True, color=WHITE, align=PP_ALIGN.LEFT
)
# Authors
add_label(
    'Adarsh Singh¹   Kushal Raj Bhandari²   Jianxi Gao³   Soham Dan²   Vivek Gupta¹',
    Inches(0.3), Inches(2.55),
    Inches(29.5), Inches(0.9),
    size=Pt(34), bold=False, color=GOLD, align=PP_ALIGN.LEFT
)
# Affiliations
add_label(
    '¹Arizona State University    ²Rensselaer Polytechnic Institute    ³Microsoft',
    Inches(0.3), Inches(3.55),
    Inches(29.5), Inches(0.7),
    size=Pt(28), bold=False, color=WHITE, align=PP_ALIGN.LEFT
)

# ASU logo
add_image('asu_logo.png', Inches(30.2), Inches(0.2), Inches(2.6), Inches(1.5))
# CORAL logo
add_image('coral_logo.jpeg', Inches(33.1), Inches(0.2), Inches(2.5), Inches(2.5))

# ACL 2026 badge (gold box)
add_rect(Inches(30.0), Inches(2.1), Inches(5.7), Inches(2.7),
         fill=GOLD, line=WHITE, line_w=Pt(4))
add_label('ACL 2026', Inches(30.1), Inches(2.2), Inches(5.5), Inches(1.1),
          size=Pt(42), bold=True, color=MAROON, align=PP_ALIGN.CENTER)
add_label('Vienna, Austria', Inches(30.1), Inches(3.3), Inches(5.5), Inches(0.6),
          size=Pt(28), bold=False, color=DARK, align=PP_ALIGN.CENTER)
add_label('arXiv: 2505.14984', Inches(30.1), Inches(3.95), Inches(5.5), Inches(0.6),
          size=Pt(24), bold=False, color=MAROON, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# COLUMN 1:  Motivation | Architecture | Datasets | Metrics
# ═══════════════════════════════════════════════════════════════════════════════

y = BODY_T

# ── Motivation ────────────────────────────────────────────────────────────────
y = section_header('Why CRAFT?  Motivation', C1L, y, CW)
add_rect(C1L, y, CW, Inches(5.1), fill=WHITE, line=MID_GRAY, line_w=Pt(1))

motivation = [
    ('▸', 'Open-Domain TQA: find the right table from 169K–419K tables before answering', False),
    ('▸', 'Sparse (BM25/SPLADE): fast but misses semantic nuance', False),
    ('▸', 'Dense (DTR/BIBERT): semantic but needs fine-tuning + slow at scale', False),
    ('▸', 'No prior work combines all three signals training-free', True),
]
bullet_box(motivation, C1L + INNER, y + Inches(0.15), CW - INNER*2, Inches(4.7),
           size=Pt(23), gap=Pt(8))
y += Inches(5.1) + Inches(0.18)

# ── CRAFT Architecture ────────────────────────────────────────────────────────
y = section_header('CRAFT Architecture', C1L, y, CW)
# Architecture overview image  (2137×2153 ≈ square)
arch_h = Inches(11.4)
add_image('craft_overview.png', C1L, y, CW, arch_h)
y += arch_h + Inches(0.1)

# Stage badges strip
bw = (CW - GUTTER*2) / 3
stage_h = Inches(0.65)
stage_badge('Stage 1: SPLADE\n169K → 5,000',   C1L,              y, bw, stage_h, ORANGE)
stage_badge('Stage 2: all-mpnet\n5K → 100',     C1L + bw + GUTTER, y, bw, stage_h, BLUE)
stage_badge('Stage 3: text-emb-3\n100 → top-k', C1L + 2*(bw+GUTTER), y, bw, stage_h, GREEN)
y += stage_h + Inches(0.18)

# ── Datasets ──────────────────────────────────────────────────────────────────
y = section_header('Benchmark Datasets', C1L, y, CW)
add_rect(C1L, y, CW, Inches(3.2), fill=WHITE, line=MID_GRAY, line_w=Pt(1))
ds_lines = [
    ('■', 'NQ-Tables:  169,898 Wikipedia tables · 966 test questions (single-hop)', False),
    ('■', 'OTT-QA:     419,000 Wikipedia tables · 2,214 dev questions (multi-hop)', False),
    ('■', 'CRAFT tested zero-shot on OTT-QA — no fine-tuning', True),
]
bullet_box(ds_lines, C1L + INNER, y + Inches(0.12), CW - INNER*2, Inches(3.0),
           size=Pt(23), gap=Pt(8))
y += Inches(3.2) + Inches(0.18)

# ── Three metric callout boxes ────────────────────────────────────────────────
cbw = (CW - GUTTER*2) / 3
cch = Inches(2.0)
callout_box('86.83', 'NQ-Tables R@10', C1L,              y, cbw, cch)
callout_box('89.88', 'OTT-QA R@10',   C1L + cbw+GUTTER, y, cbw, cch)
callout_box('70%+',  'Token Reduction',C1L + 2*(cbw+GUTTER), y, cbw, cch)
y += cch + Inches(0.18)

# ═══════════════════════════════════════════════════════════════════════════════
# COLUMN 2:  Preprocessing | NQ Results | NQ Plot | Ablation | E2E
# ═══════════════════════════════════════════════════════════════════════════════

y2 = BODY_T

# ── Preprocessing Pipeline ────────────────────────────────────────────────────
y2 = section_header('Preprocessing Pipeline', C2L, y2, CW)
add_rect(C2L, y2, CW, Inches(7.0), fill=WHITE, line=MID_GRAY, line_w=Pt(1))

# Two pipeline rows as coloured boxes with arrows
def pipe_row(steps, colors, y_pos, h=Inches(0.72)):
    bcount = len(steps)
    bwidth = (CW - INNER*2 - Inches(0.35)*(bcount-1)) / bcount
    x = C2L + INNER
    for i, (label, col) in enumerate(zip(steps, colors)):
        add_rect(x, y_pos, bwidth, h, fill=col)
        add_label(label, x + Inches(0.06), y_pos + Inches(0.06),
                  bwidth - Inches(0.12), h - Inches(0.1),
                  size=Pt(19), bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER, valign='middle')
        if i < bcount - 1:
            add_label('→', x + bwidth, y_pos,
                      Inches(0.35), h,
                      size=Pt(26), bold=True, color=DARK,
                      align=PP_ALIGN.CENTER, valign='middle')
            x += bwidth + Inches(0.35)

# Row labels
add_label('Query Expansion', C2L + INNER, y2 + Inches(0.12),
          Inches(2.8), Inches(0.45), size=Pt(21), bold=True, color=MAROON)
pipe_row(
    ['Query q', 'Gemini 1.5-Flash', 'Sub-questions q₁,q₂…'],
    [BLUE,       GREEN,              GREEN],
    y2 + Inches(0.62)
)
add_label('Table Enrichment', C2L + INNER, y2 + Inches(1.55),
          Inches(2.8), Inches(0.45), size=Pt(21), bold=True, color=MAROON)
pipe_row(
    ['Raw Table', 'Gemini 1.5-Flash', 'Title + Summary', 'Mini-Table'],
    [BLUE,        GREEN,               BLUE,              RGBColor(0xE6,0x5C,0x00)],
    y2 + Inches(2.05)
)
add_label('Row Ranking → Mini-Table  (70%+ fewer tokens)',
          C2L + INNER, y2 + Inches(3.0),
          CW - INNER*2, Inches(0.5),
          size=Pt(22), bold=True, color=DARK)
add_label('Retrieval Cascade:',
          C2L + INNER, y2 + Inches(3.6),
          CW - INNER*2, Inches(0.45),
          size=Pt(21), bold=True, color=DARK)
pipe_row(
    ['SPLADE\n(sparse)', 'all-mpnet\n(dense)', 'text-emb-3\n(neural)'],
    [ORANGE,             BLUE,                 GREEN],
    y2 + Inches(4.15), h=Inches(0.82)
)
add_label('All tables     →  5,000     →  100     →  top-k',
          C2L + INNER, y2 + Inches(5.2),
          CW - INNER*2, Inches(0.5),
          size=Pt(20), bold=False, color=RGBColor(0x55,0x55,0x55),
          align=PP_ALIGN.CENTER)
y2 += Inches(7.0) + Inches(0.18)

# ── NQ-Tables Results ─────────────────────────────────────────────────────────
y2 = section_header('NQ-Tables Retrieval Results', C2L, y2, CW)
nq_data = [
    ['Model',          'Type',   'R@1',  'R@10',  'R@50'],
    ['BM25',           'Sparse', '47.84','72.90', '89.19'],
    ['SPLADE',         'Sparse', '62.11','83.33', '94.20'],
    ['BIBERT',         'Dense',  '59.40','82.25', '93.50'],
    ['DTR',            'Dense',  '51.20','75.86', '90.60'],
    ['THYME',          'Hybrid', '50.30','73.28', '88.75'],
    ['CRAFT (ours) ★', 'Hybrid', '49.84','86.83', '97.17'],
]
col_widths_nq = [Inches(4.2), Inches(1.8), Inches(1.7), Inches(1.8), Inches(1.7)]
row_h_nq = Inches(0.6)
table_shape(nq_data, col_widths_nq, row_h_nq,
            C2L, y2, best_row=5, hdr_size=Pt(21), data_size=Pt(21))
y2 += row_h_nq * len(nq_data) + Inches(0.12)

# Gold callout
add_rect(C2L, y2, CW, Inches(0.72), fill=GOLD, line=MAROON, line_w=Pt(3))
add_label('CRAFT R@10: 86.83  ·  +3.5 pp over SPLADE (83.33)  ·  No training required',
          C2L + INNER, y2 + Inches(0.08), CW - INNER*2, Inches(0.6),
          size=Pt(24), bold=True, color=MAROON, align=PP_ALIGN.CENTER)
y2 += Inches(0.72) + Inches(0.18)

# ── NQ Context-Accuracy Plot ──────────────────────────────────────────────────
y2 = section_header('Context Length vs. Accuracy (NQ-Tables)', C2L, y2, CW)
nq_plot_h = round(11.2 * 658 / 2485 * 96) / 96   # aspect correct in inches
add_image('Context_Accuracy_F1_nq_gemini (1).png',
          C2L, y2, CW, Inches(nq_plot_h))
y2 += Inches(nq_plot_h) + Inches(0.18)

# ── Ablation Studies ──────────────────────────────────────────────────────────
y2 = section_header('Ablation Studies', C2L, y2, CW)

half_w = (CW - GUTTER) / 2
# Stage contribution sub-table
stage_data = [
    ['Configuration',         'R@10'],
    ['SPLADE (Stage 1)',       '83.33'],
    ['+ Stage 2 (all-mpnet)', '82.91'],
    ['Full CRAFT (S1+S2+S3)', '86.83'],
]
table_shape(stage_data,
            [half_w * 0.68, half_w * 0.32],
            Inches(0.52), C2L, y2, best_row=2,
            hdr_size=Pt(19), data_size=Pt(19))

# Preprocessing sub-table
prep_data = [
    ['Preprocessing',       'R@10'],
    ['None (BM25)',          '72.90'],
    ['+ Table summaries',   '81.80'],
    ['+ Query expansion',   '83.30'],
    ['Full CRAFT prep.',    '86.83'],
]
table_shape(prep_data,
            [half_w * 0.68, half_w * 0.32],
            Inches(0.52), C2L + half_w + GUTTER, y2, best_row=3,
            hdr_size=Pt(19), data_size=Pt(19))

abl_rows = max(len(stage_data), len(prep_data))
y2 += Inches(0.52) * abl_rows + Inches(0.18)

# ── End-to-End QA ─────────────────────────────────────────────────────────────
y2 = section_header('End-to-End QA Performance  (NQ-Tables F1)', C2L, y2, CW)
e2e_data = [
    ['Reader Model',  'n=1',  'n=3',  'n=5'],
    ['Llama3-8B',    '39.13','39.50','40.55'],
    ['Mistral-7B',   '35.30','41.20','44.53'],
    ['Qwen2.5-7B ★','37.01','43.85','46.49'],
]
table_shape(e2e_data,
            [Inches(4.6), Inches(2.2), Inches(2.2), Inches(2.2)],
            Inches(0.6), C2L, y2, best_row=2,
            hdr_size=Pt(21), data_size=Pt(21))
y2 += Inches(0.6) * len(e2e_data) + Inches(0.1)
add_label('Mini-tables reduce tokens 70%+ → larger context → higher F1 with more tables',
          C2L + INNER, y2, CW - INNER*2, Inches(0.55),
          size=Pt(20), bold=False, color=RGBColor(0x55,0x55,0x55),
          italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# COLUMN 3:  OTT-QA | OTT-QA Plot | Token Eff | Conclusion | Contributions
# ═══════════════════════════════════════════════════════════════════════════════

y3 = BODY_T

# ── OTT-QA Results ───────────────────────────────────────────────────────────
y3 = section_header('Zero-Shot OTT-QA Results', C3L, y3, CW)
otta_data = [
    ['Model',          'Type',   'R@10',  'R@50',  'Zero-shot'],
    ['BM25',           'Sparse', '51.94', '82.10', '✓'],
    ['SPLADE',         'Sparse', '89.52', '97.50', '✓'],
    ['BIBERT',         'Dense',  '86.50', '95.40', '✗'],
    ['CRAFT (ours) ★', 'Hybrid', '89.88', '96.07', '✓'],
]
table_shape(otta_data,
            [Inches(3.6), Inches(1.7), Inches(1.7), Inches(1.7), Inches(2.5)],
            Inches(0.62), C3L, y3, best_row=3,
            hdr_size=Pt(21), data_size=Pt(21))
y3 += Inches(0.62) * len(otta_data) + Inches(0.12)

add_rect(C3L, y3, CW, Inches(0.72), fill=GOLD, line=MAROON, line_w=Pt(3))
add_label('CRAFT R@10: 89.88  ·  Best zero-shot model  ·  Beats supervised BIBERT',
          C3L + INNER, y3 + Inches(0.08), CW - INNER*2, Inches(0.6),
          size=Pt(24), bold=True, color=MAROON, align=PP_ALIGN.CENTER)
y3 += Inches(0.72) + Inches(0.18)

# ── OTT-QA Plot ───────────────────────────────────────────────────────────────
y3 = section_header('Context Length vs. Accuracy (OTT-QA)', C3L, y3, CW)
ottqa_h = round(11.2 * 762 / 2476 * 96) / 96
add_image('Context_Accuracy_F1_ottqa_large.png',
          C3L, y3, CW, Inches(ottqa_h))
y3 += Inches(ottqa_h) + Inches(0.18)

# ── Token Efficiency ─────────────────────────────────────────────────────────
y3 = section_header('Token Efficiency  (Mini-Table)', C3L, y3, CW)
tok_w = Inches(7.5)
tok_h = Inches(7.5 * 799 / 808)
tok_x = C3L + (CW - tok_w) / 2
add_image('token_consumption_table.png', tok_x, y3, tok_w, tok_h)
y3 += tok_h + Inches(0.12)
add_label('Mini-table row-ranking reduces tokens by 70%+, enabling larger LLM context windows',
          C3L + INNER, y3, CW - INNER*2, Inches(0.6),
          size=Pt(20), bold=False, color=RGBColor(0x44,0x44,0x44), italic=True)
y3 += Inches(0.6) + Inches(0.18)

# ── Conclusion ────────────────────────────────────────────────────────────────
y3 = section_header('Conclusion', C3L, y3, CW)
add_rect(C3L, y3, CW, Inches(6.0), fill=WHITE, line=MID_GRAY, line_w=Pt(1))
concl = [
    ('✔', 'First training-free cascaded pipeline for Open-Domain TQA', True),
    ('✔', 'SOTA R@10 = 86.83 on NQ-Tables — no supervised training', True),
    ('✔', 'Zero-shot R@10 = 89.88 on OTT-QA (competitive with supervised)', False),
    ('✔', '70%+ token reduction via mini-tables → better E2E F1', False),
    ('✔', 'Stage 3 recovers full performance under noisy/perturbed queries', False),
]
bullet_box(concl, C3L + INNER, y3 + Inches(0.12), CW - INNER*2, Inches(5.7),
           size=Pt(23), gap=Pt(10))
y3 += Inches(6.0) + Inches(0.18)

# ── Key Contributions ─────────────────────────────────────────────────────────
y3 = section_header('Key Contributions', C3L, y3, CW)
add_rect(C3L, y3, CW, Inches(5.6),
         fill=CREAM, line=GOLD, line_w=Pt(3))
contrib = [
    ('①', 'Training-free cascade: sparse + dense + neural — no labeled data', True),
    ('②', 'LLM preprocessing: Gemini query expansion + table summarisation', True),
    ('③', 'Mini-table representation — 70%+ token reduction', True),
    ('④', 'Comprehensive eval: NQ-Tables & OTT-QA with full ablations', True),
]
bullet_box(contrib, C3L + INNER, y3 + Inches(0.12), CW - INNER*2, Inches(5.4),
           size=Pt(23), gap=Pt(12))
y3 += Inches(5.6) + Inches(0.18)

# ── Robustness callout ────────────────────────────────────────────────────────
y3 = section_header('Robustness', C3L, y3, CW)
rob_data = [
    ['Condition',                   'R@10'],
    ['Original queries (S1+S2+S3)', '87.16'],
    ['Perturbed queries (S1+S2)',    '83.82'],
    ['Perturbed + Stage 3',          '87.16 ✓'],
]
table_shape(rob_data,
            [Inches(8.2), Inches(3.0)],
            Inches(0.58), C3L, y3, best_row=2,
            hdr_size=Pt(20), data_size=Pt(20))

# ═══════════════════════════════════════════════════════════════════════════════
# FOOTER
# ═══════════════════════════════════════════════════════════════════════════════

add_rect(0, Inches(46.0), prs.slide_width, Inches(2.0), fill=DARK)
add_label(
    '[1] Herzig et al. (2021) NAACL  |  [2] Formal et al. (2021) SIGIR  |  '
    '[3] Nguyen et al. (2022) EMNLP  |  [4] Chen et al. (2021) ICLR',
    Inches(0.3), Inches(46.1), Inches(30), Inches(0.8),
    size=Pt(20), color=MID_GRAY, align=PP_ALIGN.LEFT
)
add_label(
    'ACL 2026  ·  Vienna  ·  arXiv: 2505.14984',
    Inches(0.3), Inches(46.9), Inches(30), Inches(0.8),
    size=Pt(20), color=GOLD, bold=True, align=PP_ALIGN.LEFT
)

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUT)
size_kb = os.path.getsize(OUT) // 1024
print(f'Saved: {OUT}  ({size_kb:,} KB)')
print(f'Slide size: {prs.slide_width.inches:.0f}" × {prs.slide_height.inches:.0f}"')
