"""
CRAFT Conference Poster — v3
Modern ACL-conference style inspired by REaR poster:
  • Navy/blue primary palette + amber gold accents
  • Large readable fonts (13–14 pt body)
  • Visual zigzag pipeline with coloured stage blocks
  • Clean tables with alternating rows and gold best-result highlights
  • Rounded-corner feel via tight border radii
Poster: 36" × 48" portrait
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Palette ────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x1B, 0x3A, 0x5C)   # dark navy  (headers, title bar)
BLUE        = RGBColor(0x1E, 0x5F, 0xA6)   # medium blue (section headers)
BLUE_LIGHT  = RGBColor(0xDB, 0xEA, 0xFB)   # very light blue (body bg)
TEAL        = RGBColor(0x00, 0x7A, 0x87)   # teal accent
GOLD        = RGBColor(0xF5, 0xA6, 0x23)   # amber gold
GOLD_LIGHT  = RGBColor(0xFF, 0xF0, 0xC8)   # light gold highlight
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF7, 0xF9, 0xFC)
LIGHT_GRAY  = RGBColor(0xE8, 0xEC, 0xF0)
MID_GRAY    = RGBColor(0xC5, 0xCC, 0xD4)
DARK_TEXT   = RGBColor(0x15, 0x1E, 0x2B)
GREEN_DARK  = RGBColor(0x1B, 0x6E, 0x3A)
# Stage colours
S0_BG  = RGBColor(0xFF, 0xF5, 0xE6)   # pre-processing — peach
S0_BD  = RGBColor(0xE6, 0x8A, 0x00)   # border
S1_BG  = RGBColor(0xE8, 0xF5, 0xE9)   # stage 1 — mint green
S1_BD  = RGBColor(0x2E, 0x7D, 0x32)
S2_BG  = RGBColor(0xE3, 0xF2, 0xFD)   # stage 2 — sky blue
S2_BD  = RGBColor(0x15, 0x65, 0xC0)
S3_BG  = RGBColor(0xF3, 0xE5, 0xF5)   # stage 3 — lavender
S3_BD  = RGBColor(0x6A, 0x1B, 0x9A)
SA_BG  = RGBColor(0xE8, 0xF8, 0xF1)   # answer gen — sea-foam
SA_BD  = RGBColor(0x00, 0x69, 0x50)

# ── Dimensions ─────────────────────────────────────────────────────────────
IN      = 914400
W       = int(36 * IN)
H       = int(48 * IN)
M       = int(0.22 * IN)           # outer margin
HDR_H   = int(4.30 * IN)
BODY_Y  = HDR_H + int(0.15 * IN)
BODY_H  = H - BODY_Y - M
GAP     = int(0.20 * IN)           # column / row gap
NCOL    = 3
COL_W   = (W - 2*M - (NCOL-1)*GAP) // NCOL
COL_X   = [M + i*(COL_W + GAP) for i in range(NCOL)]
SHH     = int(0.58 * IN)           # section-header height


# ── Primitives ─────────────────────────────────────────────────────────────

def rect(slide, x, y, w, h, fill=None, border=None, bw=Pt(0.5)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.line.width = bw
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
    else:
        s.line.fill.background()
    return s

def tb(slide, x, y, w, h, text, size=Pt(13), bold=False,
       color=DARK_TEXT, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    t = slide.shapes.add_textbox(x, y, w, h)
    tf = t.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = size; r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return t

def add_para(tf, text, size=Pt(13), bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, italic=False, sb=Pt(0), sa=Pt(2)):
    p = tf.add_paragraph()
    p.alignment = align; p.space_before = sb; p.space_after = sa
    r = p.add_run(); r.text = text
    r.font.size = size; r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return p

def sec_hdr(slide, x, y, w, title, fill=BLUE, text_color=WHITE, height=None):
    """Draws a section header bar; returns (header_height, body_y)."""
    hh = height or SHH
    rect(slide, x, y, w, hh, fill=fill, bw=Pt(0))
    tb(slide, x + int(0.14*IN), y + int(0.04*IN),
       w - int(0.18*IN), hh - int(0.06*IN),
       title, size=Pt(17), bold=True, color=text_color)
    return hh, y + hh


def arrow_down(slide, cx, y, length=int(0.25*IN), color=NAVY):
    """Thin vertical arrow pointing down."""
    stem_w = int(0.06*IN)
    # Stem
    rect(slide, cx - stem_w//2, y, stem_w, length - int(0.10*IN), fill=color, bw=Pt(0))
    # Arrowhead (triangle via textbox with ▼)
    tb(slide, cx - int(0.18*IN), y + length - int(0.20*IN),
       int(0.36*IN), int(0.22*IN), "▼",
       size=Pt(10), bold=True, color=color, align=PP_ALIGN.CENTER)


def stage_badge(slide, x, y, w, h, label, num, bg, border, text_color):
    """A coloured stage block with a left-side number badge."""
    badge_w = int(0.55*IN)
    rect(slide, x, y, w, h, fill=bg, border=border, bw=Pt(1.2))
    # Badge
    rect(slide, x, y, badge_w, h, fill=border, bw=Pt(0))
    tb(slide, x, y + int(0.05*IN), badge_w, h - int(0.05*IN),
       num, size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    return badge_w   # content starts at x + badge_w


# ── Header ─────────────────────────────────────────────────────────────────

def build_header(slide):
    rect(slide, 0, 0, W, HDR_H, fill=NAVY, bw=Pt(0))
    # gold accent bar at bottom of header
    rect(slide, 0, HDR_H - int(0.08*IN), W, int(0.08*IN), fill=GOLD, bw=Pt(0))

    # Title
    t = slide.shapes.add_textbox(int(0.28*IN), int(0.22*IN),
                                  int(26.0*IN), int(2.10*IN))
    tf = t.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "CRAFT: Training-Free Cascaded Retrieval for Tabular QA"
    r.font.size = Pt(50); r.font.bold = True; r.font.color.rgb = WHITE

    # Authors
    tb(slide, int(0.28*IN), int(2.38*IN), int(26.0*IN), int(0.68*IN),
       "  *Adarsh Singh¹     *Kushal Raj Bhandari²     Jianxi Gao³     "
       "Soham Dan²     †Vivek Gupta¹",
       size=Pt(18.5), bold=False, color=GOLD)

    # Affiliations
    tb(slide, int(0.28*IN), int(3.10*IN), int(26.0*IN), int(0.50*IN),
       "  ¹Arizona State University    ²Rensselaer Polytechnic Institute    "
       "³Microsoft    *Equal contribution    †Primary supervisor",
       size=Pt(14.5), color=RGBColor(0xCC, 0xD9, 0xE8))

    # Venue box
    vx, vy = int(26.7*IN), int(0.22*IN)
    vw, vh = int(5.10*IN), int(1.18*IN)
    rect(slide, vx, vy, vw, vh, fill=BLUE, border=GOLD, bw=Pt(2.5))
    tb(slide, vx + int(0.12*IN), vy + int(0.10*IN), vw - int(0.18*IN), vh,
       "ACL 2026  ·  Vienna, Austria",
       size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # arXiv
    tb(slide, vx + int(0.12*IN), vy + vh + int(0.12*IN),
       vw - int(0.18*IN), int(0.45*IN),
       "arXiv: 2505.14984",
       size=Pt(15.5), bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Emails
    tb(slide, vx + int(0.12*IN), vy + vh + int(0.65*IN),
       vw - int(0.18*IN), int(0.80*IN),
       "asing725@asu.edu\nbhandk@rpi.edu",
       size=Pt(13), color=RGBColor(0xCC, 0xD9, 0xE8), align=PP_ALIGN.CENTER)

    # Logos
    for path, lx, lh in [
        ("/mnt/data1/asing725/ACL/CRAFT/static/images/asu_logo.png",
         int(32.1*IN), int(0.75*IN)),
        ("/mnt/data1/asing725/ACL/CRAFT/static/images/coral_logo.jpeg",
         int(34.1*IN), int(0.85*IN)),
    ]:
        try:
            img = slide.shapes.add_picture(path, lx, int(3.25*IN), int(1.75*IN), lh)
        except Exception:
            pass


# ── Section: Motivation ────────────────────────────────────────────────────

def sec_motivation(slide, x, y, w, h):
    hh, by = sec_hdr(slide, x, y, w, "  Motivation")
    rect(slide, x, by, w, h - hh, fill=OFF_WHITE, border=MID_GRAY, bw=Pt(0.5))
    bx = x + int(0.14*IN); bw2 = w - int(0.28*IN)

    t = slide.shapes.add_textbox(bx, by + int(0.12*IN), bw2, h - hh - int(0.14*IN))
    tf = t.text_frame; tf.word_wrap = True

    paras = [
        ("Open-Domain Table QA", True, Pt(14.5), NAVY),
        ("Retrieve the right table from a large corpus of 169K–419K tables, "
         "then answer natural language queries.", False, Pt(13.5), DARK_TEXT),
        ("", False, Pt(4), DARK_TEXT),
        ("Motivating Example", True, Pt(14.5), NAVY),
        ('Query: "Find Customers who bought products from suppliers in Japan"', False, Pt(13), TEAL),
        ("  Retrieved (Top-2):  Products,  Suppliers  ✗  No join path → SQL fails!", False, Pt(13), DARK_TEXT),
        ('  Missing but needed:  Customers  ✓', False, Pt(13), GREEN_DARK),
        ("", False, Pt(4), DARK_TEXT),
        ("The Problem", True, Pt(14.5), NAVY),
        ("▸  Sparse / Dense: optimise query–table relevance only; ignore join-ability", False, Pt(13), DARK_TEXT),
        ("▸  LLM-based (ARM, JAR): expensive loops → high latency & token cost", False, Pt(13), DARK_TEXT),
        ("▸  Larger top-k: more recall but noisy, weakly related tables hurt precision", False, Pt(13), DARK_TEXT),
        ("", False, Pt(4), DARK_TEXT),
        ("Research Question", True, Pt(14.5), NAVY),
        ("Can off-the-shelf pretrained models, combined in a zero-shot cascade, "
         "match costly fine-tuned retrievers at scale?", False, Pt(13.5), RGBColor(0x15,0x65,0xC0)),
    ]
    first = True
    for txt, bold, sz, clr in paras:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(1)
        r = p.add_run(); r.text = txt
        r.font.size = sz; r.font.bold = bold; r.font.color.rgb = clr


# ── Section: Why Existing Fails ────────────────────────────────────────────

def sec_why_fails(slide, x, y, w, h):
    hh, by = sec_hdr(slide, x, y, w, "  Why Existing Approaches Fail")
    rect(slide, x, by, w, h - hh, fill=OFF_WHITE, border=MID_GRAY, bw=Pt(0.5))

    rows = [
        ("Approach", "Problem", True),
        ("Dense / Sparse\nRetrieval", "Optimise semantic relevance only;\nignore joinability & connectivity", False),
        ("LLM-based\n(ARM, JAR)", "Expensive LLM loops during retrieval\n→ high latency and token cost", False),
        ("Larger\nTop-k", "Improves recall but noisy, weakly related\ntables hurt downstream precision", False),
    ]
    bx = x + int(0.10*IN); bw2 = w - int(0.20*IN)
    row_h = int((h - hh - int(0.14*IN)) / len(rows))
    col1 = int(bw2 * 0.37); col2 = bw2 - col1
    ty = by + int(0.07*IN)
    fills = [NAVY, LIGHT_GRAY, BLUE_LIGHT, RGBColor(0xFF,0xF5,0xE6)]
    text_colors = [WHITE, DARK_TEXT, DARK_TEXT, DARK_TEXT]

    for ri, (a, b, is_hdr) in enumerate(rows):
        fill = fills[ri]
        tc = text_colors[ri]
        # col A
        rect(slide, bx, ty, col1, row_h, fill=fill, border=MID_GRAY, bw=Pt(0.4))
        tb(slide, bx + int(0.07*IN), ty + int(0.04*IN),
           col1 - int(0.10*IN), row_h - int(0.06*IN),
           a, size=Pt(13 if not is_hdr else 13.5), bold=is_hdr, color=tc)
        # col B
        rect(slide, bx + col1, ty, col2, row_h, fill=fill, border=MID_GRAY, bw=Pt(0.4))
        tb(slide, bx + col1 + int(0.07*IN), ty + int(0.04*IN),
           col2 - int(0.10*IN), row_h - int(0.06*IN),
           b, size=Pt(12.5 if not is_hdr else 13.5), bold=is_hdr, color=tc)
        ty += row_h


# ── Section: Key Insight ───────────────────────────────────────────────────

def sec_key_insight(slide, x, y, w, h):
    hh, by = sec_hdr(slide, x, y, w, "  Key Insight", fill=RGBColor(0xB7,0x7D,0x00),
                     text_color=WHITE)
    rect(slide, x, by, w, h - hh, fill=GOLD_LIGHT, border=GOLD, bw=Pt(1.2))
    bx = x + int(0.14*IN); bw2 = w - int(0.28*IN)

    t = slide.shapes.add_textbox(bx, by + int(0.10*IN), bw2, h - hh - int(0.14*IN))
    tf = t.text_frame; tf.word_wrap = True

    paras = [
        ("CRAFT jointly models:", True, Pt(14), NAVY),
        ("", False, Pt(3), DARK_TEXT),
        ("  Query ↔ Table  Relevance", True, Pt(14), BLUE),
        ("  (What is relevant?)", False, Pt(13), DARK_TEXT),
        ("", False, Pt(2), DARK_TEXT),
        ("  Table ↔ Table  Joinability", True, Pt(14), RGBColor(0x1B,0x6E,0x3A)),
        ("  (Can they be connected?)", False, Pt(13), DARK_TEXT),
        ("", False, Pt(2), DARK_TEXT),
        ("  ⟹  Join-ready, high-fidelity table sets", True, Pt(14), NAVY),
        ("       without LLM calls", True, Pt(14), NAVY),
    ]
    first = True
    for txt, bold, sz, clr in paras:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(1)
        r = p.add_run(); r.text = txt
        r.font.size = sz; r.font.bold = bold; r.font.color.rgb = clr


# ── Section: Contributions ─────────────────────────────────────────────────

def sec_contributions(slide, x, y, w, h):
    hh, by = sec_hdr(slide, x, y, w, "  Key Contributions")
    rect(slide, x, by, w, h - hh, fill=OFF_WHITE, border=MID_GRAY, bw=Pt(0.5))
    bx = x + int(0.14*IN); bw2 = w - int(0.28*IN)

    t = slide.shapes.add_textbox(bx, by + int(0.12*IN), bw2, h - hh - int(0.15*IN))
    tf = t.text_frame; tf.word_wrap = True
    items = [
        ("① Training-Free Zero-Shot Cascade",
         "No dataset-specific fine-tuning; plug-and-play on any domain."),
        ("② New SOTA on NQ-Tables",
         "Recall@1 = 49.84, R@10 = 86.83, R@50 = 97.17 — no training required."),
        ("③ Robust to Query Paraphrasing",
         "Only –0.04 Δ on perturbed queries vs. –5 to –12 for DTR variants."),
        ("④ 33× Fewer Embedding Calls",
         "Mini-tables cut >70% tokens; offline pre-encoding avoids online LLM use."),
    ]
    first = True
    for title, body in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph(); p.space_before = Pt(5)
        r = p.add_run(); r.text = title
        r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = NAVY
        p2 = tf.add_paragraph(); p2.space_after = Pt(1)
        r2 = p2.add_run(); r2.text = body
        r2.font.size = Pt(13); r2.font.color.rgb = DARK_TEXT


# ── Section: CRAFT Overview figure ─────────────────────────────────────────

def sec_overview_fig(slide, x, y, w, h):
    hh, by = sec_hdr(slide, x, y, w, "  CRAFT Framework Overview")
    rect(slide, x, by, w, h - hh, fill=WHITE, border=MID_GRAY, bw=Pt(0.5))
    try:
        pad = int(0.12*IN)
        slide.shapes.add_picture(
            "/mnt/data1/asing725/ACL/CRAFT/static/images/craft_overview.png",
            x + pad, by + pad, w - 2*pad, h - hh - 2*pad)
    except Exception as e:
        tb(slide, x + int(0.15*IN), by + int(0.2*IN), w - int(0.3*IN), int(1.0*IN),
           f"[craft_overview.png — {e}]", size=Pt(11), color=DARK_TEXT)


# ── Section: Pipeline (zigzag) ─────────────────────────────────────────────

def sec_pipeline(slide, x, y, w, h):
    hh, by = sec_hdr(slide, x, y, w, "  CRAFT Pipeline  ·  Three-Stage Cascade")
    rect(slide, x, by, w, h - hh, fill=OFF_WHITE, border=MID_GRAY, bw=Pt(0.5))

    bx = x + int(0.12*IN)
    bw = w - int(0.24*IN)
    cy = by + int(0.12*IN)

    arrow_h = int(0.32*IN)

    # ── Pre-processing block ───────────────────────────────────────────────
    pre_h = int(2.05*IN)
    rect(slide, bx, cy, bw, pre_h, fill=S0_BG, border=S0_BD, bw=Pt(1.5))
    # label stripe
    label_w = int(1.05*IN)
    rect(slide, bx, cy, label_w, pre_h, fill=S0_BD, bw=Pt(0))
    # rotated label via textbox
    tb(slide, bx + int(0.04*IN), cy + int(0.04*IN), label_w - int(0.06*IN), pre_h - int(0.06*IN),
       "PRE-\nPROC.", size=Pt(12.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    cx = bx + label_w + int(0.10*IN)
    cw = bw - label_w - int(0.14*IN)
    t = slide.shapes.add_textbox(cx, cy + int(0.08*IN), cw, pre_h - int(0.12*IN))
    tf = t.text_frame; tf.word_wrap = True

    def _pline(txt, bold=False, sz=Pt(13), clr=DARK_TEXT):
        p = tf.add_paragraph(); p.space_after = Pt(1)
        r = p.add_run(); r.text = txt
        r.font.size = sz; r.font.bold = bold; r.font.color.rgb = clr

    p0 = tf.paragraphs[0]; p0.space_after = Pt(1)
    r0 = p0.add_run(); r0.text = "Gemini Flash 1.5  ·  Applied once offline"
    r0.font.size = Pt(14); r0.font.bold = True; r0.font.color.rgb = RGBColor(0xB7,0x5A,0x00)
    _pline("▸  Query Decomposition — break queries into sub-questions")
    _pline("▸  Table Enrichment — generate descriptive titles + summaries")
    _pline("▸  Row Ranking — rank rows by semantic relevance (all-mpnet-base-v2)")
    cy += pre_h

    # ── Stage blocks ──────────────────────────────────────────────────────
    stage_cfgs = [
        ("STAGE\n  1", "Sparse Lexical Retrieval  ·  SPLADE",
         "▸  Sparse lexical-expansion over full corpus\n"
         "   (169K NQ-Tables / 419K OTT-QA)\n"
         "▸  Enriched table descriptions + sub-questions\n"
         "   boost first-stage recall significantly",
         "→  N₁ = 5,000\n    candidates",
         S1_BG, S1_BD, False),
        ("STAGE\n  2", "Dense Semantic Pruning  ·  Sentence Transformer",
         "▸  all-mpnet-base-v2 (NQ) / Jina Embeddings v3 (OTT-QA)\n"
         "▸  Mini-tables (top-5 rows) scored against query\n"
         "▸  Sublinear ANN search over pre-encoded row embeddings",
         "→  N₂ = 100–300\n    mini-tables",
         S2_BG, S2_BD, True),
        ("STAGE\n  3", "Neural Re-ranking  ·  text-embedding-3 / gemini-embed",
         "▸  text-embedding-3-small / large  (NQ-Tables)\n"
         "▸  gemini-embedding-001  (OTT-QA)\n"
         "▸  Precision-optimised final top-K selection",
         "→  Top-K tables\n    (K ∈ 1, 3, 5…10)",
         S3_BG, S3_BD, False),
    ]

    stage_h = int(2.60*IN)
    output_w = int(1.80*IN)
    stage_bw = bw - output_w - int(0.10*IN)

    for label, title, body, output, bg, bd, flip in stage_cfgs:
        # arrow down
        arrow_cx = bx + bw // 2
        rect(slide, arrow_cx - int(0.04*IN), cy, int(0.08*IN),
             arrow_h - int(0.10*IN), fill=bd, bw=Pt(0))
        tb(slide, arrow_cx - int(0.15*IN), cy + arrow_h - int(0.20*IN),
           int(0.30*IN), int(0.22*IN), "▼",
           size=Pt(10), bold=True, color=bd, align=PP_ALIGN.CENTER)
        cy += arrow_h

        # Stage block
        sx = bx
        rect(slide, sx, cy, stage_bw, stage_h, fill=bg, border=bd, bw=Pt(1.5))
        # badge
        badge_w = int(0.60*IN)
        rect(slide, sx, cy, badge_w, stage_h, fill=bd, bw=Pt(0))
        tb(slide, sx + int(0.02*IN), cy + int(0.05*IN),
           badge_w - int(0.04*IN), stage_h - int(0.08*IN),
           label, size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Content
        inner_x = sx + badge_w + int(0.10*IN)
        inner_w = stage_bw - badge_w - int(0.14*IN)
        tb(slide, inner_x, cy + int(0.06*IN), inner_w, int(0.42*IN),
           title, size=Pt(13.5), bold=True, color=bd)
        tb(slide, inner_x, cy + int(0.48*IN), inner_w,
           stage_h - int(0.52*IN),
           body, size=Pt(12.5), color=DARK_TEXT)

        # Output bubble (right side)
        ob_x = sx + stage_bw + int(0.10*IN)
        ob_y = cy + (stage_h - int(0.95*IN)) // 2
        rect(slide, ob_x, ob_y, output_w - int(0.12*IN), int(0.95*IN),
             fill=bd, bw=Pt(0))
        tb(slide, ob_x + int(0.06*IN), ob_y + int(0.04*IN),
           output_w - int(0.20*IN), int(0.90*IN),
           output, size=Pt(12), bold=True, color=WHITE)

        cy += stage_h

    # Answer generation
    rect(slide, bx, cy + int(0.32*IN), bw, int(1.25*IN),
         fill=SA_BG, border=SA_BD, bw=Pt(1.5))
    rect(slide, bx, cy + int(0.32*IN), int(0.60*IN), int(1.25*IN), fill=SA_BD, bw=Pt(0))
    tb(slide, bx + int(0.02*IN), cy + int(0.34*IN),
       int(0.58*IN), int(1.22*IN),
       "ANS\nGEN", size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    rect(slide, bx + int(0.60*IN), cy + int(0.32*IN),
         int(0.06*IN), int(1.25*IN), fill=SA_BD, bw=Pt(0))

    tb(slide, bx + int(0.80*IN), cy + int(0.40*IN),
       bw - int(0.90*IN), int(1.10*IN),
       "Pretrained LLM (Llama3-8B / Qwen2.5-7B / Mistral-7B)\n"
       "Top-K mini-tables  +  query  →  Final Answer",
       size=Pt(13), bold=False, color=DARK_TEXT)

    # Down-arrow before answer gen
    mid_x = bx + bw // 2
    rect(slide, mid_x - int(0.04*IN), cy, int(0.08*IN),
         int(0.25*IN), fill=SA_BD, bw=Pt(0))
    tb(slide, mid_x - int(0.15*IN), cy + int(0.22*IN),
       int(0.30*IN), int(0.22*IN), "▼",
       size=Pt(10), bold=True, color=SA_BD, align=PP_ALIGN.CENTER)


# ── Shared table builder ───────────────────────────────────────────────────

def draw_table(slide, bx, ty, bw, col_fracs, headers, data_rows,
               row_h=int(0.30*IN), cat_groups=None, hdr_fill=NAVY):
    """
    data_rows: list of (cells..., category_key, is_bold)
    cat_groups: dict {key: fill_color}
    """
    if cat_groups is None:
        cat_groups = {}
    col_ws = [int(bw * f) for f in col_fracs]

    def _cell(cx, ry, cw, ch, text, fill, text_c, bold, size, align):
        rect(slide, cx, ry, cw, ch, fill=fill, border=MID_GRAY, bw=Pt(0.4))
        tb(slide, cx + int(0.04*IN), ry + int(0.02*IN),
           cw - int(0.06*IN), ch - int(0.03*IN),
           text, size=size, bold=bold, color=text_c, align=align)

    # Header row
    cx = bx
    for ci, (hdr, cw) in enumerate(zip(headers, col_ws)):
        _cell(cx, ty, cw, row_h, hdr, hdr_fill, WHITE, True, Pt(12.5),
              PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
        cx += cw
    ty += row_h

    # Data rows
    default_fills = [LIGHT_GRAY, OFF_WHITE] * 30
    for ri, row in enumerate(data_rows):
        *cells, cat, bold = row
        fill = cat_groups.get(cat, default_fills[ri % 2])
        text_c = NAVY if (cat == "craft" or cat == "cat") else DARK_TEXT
        cx = bx
        for ci, (val, cw) in enumerate(zip(cells, col_ws)):
            _cell(cx, ty, cw, row_h, val, fill, text_c, bold,
                  Pt(12 if cat != "cat" else 11),
                  PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
            cx += cw
        ty += row_h
    return ty


CAT_FILLS = {
    "cat":     RGBColor(0xC8, 0xD8, 0xEC),  # category header — muted blue
    "sparse":  OFF_WHITE,
    "dense":   BLUE_LIGHT,
    "hybrid":  RGBColor(0xE8, 0xF5, 0xE9),
    "craft_s": RGBColor(0xFF, 0xF5, 0xE0),
    "craft":   GOLD_LIGHT,
    "base":    OFF_WHITE,
    "bad":     RGBColor(0xFF, 0xED, 0xED),
}


# ── Section: NQ-Tables Retrieval ───────────────────────────────────────────

def sec_nq(slide, x, y, w, h):
    hh, by = sec_hdr(slide, x, y, w, "  Retrieval  ·  NQ-Tables (169K tables)")
    rect(slide, x, by, w, h - hh, fill=OFF_WHITE, border=MID_GRAY, bw=Pt(0.5))
    bx = x + int(0.10*IN); bw = w - int(0.20*IN)

    tb(slide, bx, by + int(0.08*IN), bw, int(0.55*IN),
       "All models except CRAFT are trained on NQ-Tables. R = Recall. "
       "Bold = best; underlined = second-best.",
       size=Pt(12), color=DARK_TEXT, italic=True)

    ty = by + int(0.65*IN)
    rows = [
        ("Sparse Retrieval", "", "", "", "cat", True),
        ("BM25",             "18.49", "36.94", "52.61", "sparse", False),
        ("SPLADE",           "39.84", "83.33", "94.65", "sparse", False),
        ("Dense Retrieval",  "", "", "", "cat", True),
        ("BIBERT",           "43.78", "82.25", "93.71", "dense",  False),
        ("DPR",              "45.32", "85.84", "95.44", "dense",  False),
        ("DTR",              "32.62", "75.86", "89.77", "dense",  False),
        ("SSDR",             "45.47", "84.00", "95.05", "dense",  False),
        ("Hybrid Retrieval", "", "", "", "cat", True),
        ("BIBERT+SPLADE",    "43.67", "86.72", "95.62", "hybrid", False),
        ("THYME",            "48.55", "86.38", "96.08", "hybrid", False),
        ("CRAFT Stage-wise", "", "", "", "cat", True),
        ("Stage 1",          "34.38", "72.90", "91.62", "craft_s", False),
        ("+ Stage 2",        "36.65", "82.91", "96.08", "craft_s", False),
        ("+ Stage 3",        "41.13", "87.16", "96.84", "craft_s", False),
        ("CRAFT (full) ★",   "49.84", "86.83", "97.17", "craft",  True),
    ]
    row_h = min(int((h - hh - int(0.75*IN) - int(0.70*IN)) / len(rows)), int(0.29*IN))
    ty = draw_table(slide, bx, ty, bw, [0.44, 0.19, 0.19, 0.18],
                    ["Model", "R@1", "R@10", "R@50"], rows,
                    row_h=row_h, cat_groups=CAT_FILLS)

    # Callout
    callout_y = ty + int(0.08*IN)
    available = (y + h) - callout_y - int(0.05*IN)
    if available > int(0.40*IN):
        ch = min(int(0.58*IN), available)
        rect(slide, bx, callout_y, bw, ch, fill=GOLD_LIGHT, border=GOLD, bw=Pt(1.5))
        tb(slide, bx + int(0.10*IN), callout_y + int(0.05*IN),
           bw - int(0.15*IN), ch - int(0.06*IN),
           "★  CRAFT: new SOTA  R@1 = 49.84  on NQ-Tables  (no training)",
           size=Pt(13.5), bold=True, color=NAVY)


# ── Section: OTT-QA Retrieval ──────────────────────────────────────────────

def sec_ottqa(slide, x, y, w, h):
    hh, by = sec_hdr(slide, x, y, w, "  Retrieval  ·  OTT-QA (Zero-Shot, 419K tables)")
    rect(slide, x, by, w, h - hh, fill=OFF_WHITE, border=MID_GRAY, bw=Pt(0.5))
    bx = x + int(0.10*IN); bw = w - int(0.20*IN)

    tb(slide, bx, by + int(0.08*IN), bw, int(0.55*IN),
       "OTT-QA requires multi-hop reasoning over tables + passages. "
       "CRAFT is fully zero-shot — all other methods are fine-tuned.",
       size=Pt(12), color=DARK_TEXT, italic=True)

    ty = by + int(0.65*IN)
    rows = [
        ("Sparse Retrieval", "", "", "", "cat", True),
        ("BM25",             "23.98", "51.94", "69.11", "sparse", False),
        ("SPLADE",           "62.74", "89.52", "95.21", "sparse", False),
        ("Dense Retrieval",  "", "", "", "cat", True),
        ("BIBERT",           "56.82", "86.50", "94.26", "dense",  False),
        ("DPR",              "53.43", "85.95", "93.22", "dense",  False),
        ("SSDR",             "56.96", "86.22", "93.55", "dense",  False),
        ("Hybrid Retrieval", "", "", "", "cat", True),
        ("BIBERT+SPLADE",    "64.72", "91.01", "96.34", "hybrid", False),
        ("THYME ★",          "66.67", "91.10", "96.16", "hybrid", False),
        ("CRAFT Stage-wise", "", "", "", "cat", True),
        ("Stage 1",          "40.74", "60.89", "88.17", "craft_s", False),
        ("+ Stage 2",        "47.33", "87.35", "91.82", "craft_s", False),
        ("CRAFT (full)",     "55.56", "89.88", "96.07", "craft",  True),
    ]
    row_h = min(int((h - hh - int(0.75*IN) - int(0.70*IN)) / len(rows)), int(0.29*IN))
    ty = draw_table(slide, bx, ty, bw, [0.44, 0.19, 0.19, 0.18],
                    ["Model", "R@1", "R@10", "R@50"], rows,
                    row_h=row_h, cat_groups=CAT_FILLS)

    callout_y = ty + int(0.08*IN)
    available = (y + h) - callout_y - int(0.05*IN)
    if available > int(0.38*IN):
        ch = min(int(0.60*IN), available)
        rect(slide, bx, callout_y, bw, ch, fill=BLUE_LIGHT, border=BLUE, bw=Pt(1.2))
        tb(slide, bx + int(0.10*IN), callout_y + int(0.05*IN),
           bw - int(0.15*IN), ch - int(0.06*IN),
           "Zero-shot CRAFT approaches fine-tuned THYME at R@10 & R@50",
           size=Pt(13), bold=True, color=NAVY)


# ── Section: End-to-End QA ─────────────────────────────────────────────────

def sec_e2e(slide, x, y, w, h):
    hh, by = sec_hdr(slide, x, y, w,
                     "  End-to-End QA (F1)  ·  NQ-Tables")
    rect(slide, x, by, w, h - hh, fill=OFF_WHITE, border=MID_GRAY, bw=Pt(0.5))
    bx = x + int(0.10*IN); bw = w - int(0.20*IN)

    tb(slide, bx, by + int(0.08*IN), bw, int(0.55*IN),
       "F1 with varying retrieved tables n. Models: Mistral-7B, Llama3-8B, Qwen2.5-7B.",
       size=Pt(12), color=DARK_TEXT, italic=True)

    ty = by + int(0.65*IN)
    row_h = int(0.30*IN)
    col_fracs = [0.26, 0.083, 0.083, 0.083, 0.083, 0.083, 0.083, 0.083, 0.083, 0.083]
    col_ws = [int(bw * f) for f in col_fracs]

    # Span headers: n=1, n=3, n=5
    span_labels = [("", 1, NAVY), ("n = 1", 3, BLUE),
                   ("n = 3", 3, RGBColor(0x1B,0x6E,0x3A)),
                   ("n = 5", 3, RGBColor(0x6A,0x1B,0x9A))]
    cx = bx
    for label, span, color in span_labels:
        sw = sum(col_ws[sum(s for _,s,_ in span_labels[:span_labels.index((label,span,color))]):
                        sum(s for _,s,_ in span_labels[:span_labels.index((label,span,color))]) + span])
        rect(slide, cx, ty, sw, row_h, fill=color, bw=Pt(0))
        tb(slide, cx, ty, sw, row_h, label, size=Pt(12), bold=True,
           color=WHITE, align=PP_ALIGN.CENTER)
        cx += sw
    ty += row_h

    # Sub-header
    sub = ["Retriever", "Mis.", "Ll3", "Qwn", "Mis.", "Ll3", "Qwn", "Mis.", "Ll3", "Qwn"]
    sub_colors = [NAVY, BLUE, BLUE, BLUE,
                  RGBColor(0x1B,0x6E,0x3A), RGBColor(0x1B,0x6E,0x3A), RGBColor(0x1B,0x6E,0x3A),
                  RGBColor(0x6A,0x1B,0x9A), RGBColor(0x6A,0x1B,0x9A), RGBColor(0x6A,0x1B,0x9A)]
    cx = bx
    for sh, cw, sc in zip(sub, col_ws, sub_colors):
        rect(slide, cx, ty, cw, row_h, fill=RGBColor(0xCC, 0xD9, 0xEC), border=MID_GRAY, bw=Pt(0.3))
        tb(slide, cx, ty, cw, row_h, sh, size=Pt(11.5), bold=True, color=sc, align=PP_ALIGN.CENTER)
        cx += cw
    ty += row_h

    data = [
        ("BIBERT",        "32.93","32.66","34.80","34.6","33.16","37.92","35.30","34.28","37.01","base",False),
        ("SPLADE",        "29.61","32.07","31.90","35.42","37.17","35.79","34.95","33.88","37.35","base",False),
        ("BIBERT+SPLADE", "32.67","32.66","33.24","33.59","33.66","36.76","35.71","33.92","37.02","base",False),
        ("THYME",         "35.48","36.14","37.28","37.59","39.16","40.28","37.20","39.29","41.20","hybrid",False),
        ("CRAFT ★",       "39.13","38.31","39.73","45.28","40.76","43.52","44.53","40.55","46.49","craft",True),
    ]
    row_fills = {"base": OFF_WHITE, "hybrid": RGBColor(0xE8,0xF5,0xE9), "craft": GOLD_LIGHT}

    for ri, row in enumerate(data):
        *cells, cat, bold = row
        fill = row_fills.get(cat, OFF_WHITE)
        text_c = NAVY if bold else DARK_TEXT
        cx = bx
        for ci, (val, cw) in enumerate(zip(cells, col_ws)):
            rect(slide, cx, ty, cw, row_h, fill=fill, border=MID_GRAY, bw=Pt(0.3))
            tb(slide, cx + int(0.03*IN), ty + int(0.02*IN), cw - int(0.05*IN),
               row_h - int(0.03*IN), val,
               size=Pt(12 if ci == 0 else 11.5), bold=bold, color=text_c,
               align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
            cx += cw
        ty += row_h

    # callout
    callout_y = ty + int(0.08*IN)
    available = (y + h) - callout_y - int(0.05*IN)
    if available > int(0.38*IN):
        ch = min(int(0.55*IN), available)
        rect(slide, bx, callout_y, bw, ch, fill=GOLD_LIGHT, border=GOLD, bw=Pt(1.5))
        tb(slide, bx + int(0.10*IN), callout_y + int(0.05*IN),
           bw - int(0.15*IN), ch - int(0.07*IN),
           "★  CRAFT tops F1 at n=1 and n=5 across all three LLMs",
           size=Pt(13.5), bold=True, color=NAVY)


# ── Section: Robustness ────────────────────────────────────────────────────

def sec_robustness(slide, x, y, w, h):
    hh, by = sec_hdr(slide, x, y, w, "  Robustness to Query Paraphrasing",
                     fill=RGBColor(0x15, 0x65, 0xC0))
    rect(slide, x, by, w, h - hh, fill=OFF_WHITE, border=MID_GRAY, bw=Pt(0.5))
    bx = x + int(0.10*IN); bw = w - int(0.20*IN)

    tb(slide, bx, by + int(0.08*IN), bw, int(0.50*IN),
       "Perturbed queries via Gemini 2.5 Flash. Δ = avg change from original query.",
       size=Pt(12), italic=True, color=DARK_TEXT)

    rows = [
        ("", "R@1", "R@10", "R@50", "Δ Avg", "cat", True),
        ("DTR (M) — original",  "38.74","75.73","88.36","—", "base", False),
        ("CRAFT — original",    "41.13","87.16","96.84","—", "craft",True),
        ("DTR (M) — perturbed", "27.09","66.27","84.33","–8.38","bad",False),
        ("DTR+HN — perturbed",  "37.87","76.06","84.43","–5.80","bad",False),
        ("CRAFT — perturbed",   "41.02","86.83","96.08","–0.04","craft",True),
    ]
    cat_fills_rob = {**CAT_FILLS, "bad": RGBColor(0xFF,0xED,0xED), "base": BLUE_LIGHT}
    row_h = int(0.30*IN)
    ty = draw_table(slide, bx, by + int(0.60*IN), bw,
                    [0.37,0.15,0.15,0.15,0.18],
                    ["Model","R@1","R@10","R@50","Δ Avg"],
                    rows[1:], row_h=row_h, cat_groups=cat_fills_rob)

    callout_y = ty + int(0.08*IN)
    available = (y + h) - callout_y - int(0.05*IN)
    if available > int(0.38*IN):
        ch = min(int(0.58*IN), available)
        rect(slide, bx, callout_y, bw, ch,
             fill=BLUE_LIGHT, border=BLUE, bw=Pt(1.2))
        tb(slide, bx + int(0.10*IN), callout_y + int(0.06*IN),
           bw - int(0.15*IN), ch - int(0.07*IN),
           "CRAFT drops only –0.04 on paraphrased queries vs. –5.8 to –8.4 for DTR",
           size=Pt(13), bold=True, color=NAVY)


# ── Section: Stage-wise Ablation (bar chart) ───────────────────────────────

def sec_ablation(slide, x, y, w, h):
    hh, by = sec_hdr(slide, x, y, w, "  Ablation  ·  Stage-wise Gains (R@10)",
                     fill=RGBColor(0x5C, 0x1A, 0x5C))
    rect(slide, x, by, w, h - hh, fill=OFF_WHITE, border=MID_GRAY, bw=Pt(0.5))
    bx = x + int(0.12*IN); bw = w - int(0.24*IN)

    tb(slide, bx, by + int(0.08*IN), bw, int(0.42*IN),
       "NQ-Tables Recall@10 at each stage of the cascade:",
       size=Pt(13), bold=True, color=DARK_TEXT)

    bars = [
        ("Stage 1\nSPLADE",     72.90, S1_BD),
        ("+ Stage 2\nDense",    82.91, S2_BD),
        ("+ Stage 3\nRerank",   87.16, S3_BD),
        ("CRAFT\n(full)",       86.83, NAVY),
    ]

    chart_y = by + int(0.52*IN)
    chart_h = h - hh - int(0.55*IN) - int(0.50*IN)
    axis_h  = chart_h - int(0.45*IN)   # bar area
    max_v   = 100.0

    bar_total_w = bw - int(0.40*IN)
    bar_w  = int(bar_total_w / len(bars)) - int(0.14*IN)
    bar_gx = int(0.40*IN)

    for i, (label, val, color) in enumerate(bars):
        bxi = bx + bar_gx + i * (bar_w + int(0.14*IN))
        filled_h = int((val / max_v) * axis_h)
        bar_top  = chart_y + axis_h - filled_h

        # Background track
        rect(slide, bxi, chart_y, bar_w, axis_h,
             fill=RGBColor(0xE8,0xEC,0xF0), bw=Pt(0))
        # Filled bar
        rect(slide, bxi, bar_top, bar_w, filled_h, fill=color, bw=Pt(0))
        # Value
        tb(slide, bxi - int(0.06*IN), bar_top - int(0.32*IN),
           bar_w + int(0.12*IN), int(0.30*IN),
           f"{val:.1f}", size=Pt(13), bold=True, color=color, align=PP_ALIGN.CENTER)
        # X label
        tb(slide, bxi - int(0.04*IN), chart_y + axis_h + int(0.04*IN),
           bar_w + int(0.08*IN), int(0.42*IN),
           label, size=Pt(11.5), color=DARK_TEXT, align=PP_ALIGN.CENTER)

    # Insight box at bottom
    insight_y = chart_y + chart_h + int(0.06*IN)
    rem = (y + h) - insight_y - int(0.05*IN)
    if rem > int(0.35*IN):
        ch = min(int(0.52*IN), rem)
        rect(slide, bx, insight_y, bw, ch,
             fill=RGBColor(0xF3,0xE5,0xF5), border=S3_BD, bw=Pt(1.0))
        tb(slide, bx + int(0.08*IN), insight_y + int(0.05*IN),
           bw - int(0.12*IN), ch - int(0.06*IN),
           "Each stage contributes distinct, non-overlapping recall gains",
           size=Pt(13), bold=True, color=S3_BD)


# ── Section: Efficiency ────────────────────────────────────────────────────

def sec_efficiency(slide, x, y, w, h):
    hh, by = sec_hdr(slide, x, y, w, "  Efficiency  ·  Mini-Tables & Embedding Calls",
                     fill=GREEN_DARK)
    rect(slide, x, by, w, h - hh, fill=OFF_WHITE, border=MID_GRAY, bw=Pt(0.5))
    bx = x + int(0.12*IN); bw = w - int(0.24*IN)

    # Three stat boxes
    stats = [
        ("33×", "fewer\nembedding calls", S1_BD),
        (">70%", "token\nreduction", BLUE),
        ("0", "online\nLLM calls", RGBColor(0x6A,0x1B,0x9A)),
    ]
    stat_bw = (bw - 2*int(0.10*IN)) // 3
    sx = bx; sy = by + int(0.12*IN)
    for big, small, color in stats:
        rect(slide, sx, sy, stat_bw, int(1.10*IN), fill=WHITE, border=color, bw=Pt(1.5))
        tb(slide, sx, sy + int(0.04*IN), stat_bw, int(0.55*IN),
           big, size=Pt(28), bold=True, color=color, align=PP_ALIGN.CENTER)
        tb(slide, sx + int(0.04*IN), sy + int(0.56*IN), stat_bw - int(0.06*IN), int(0.52*IN),
           small, size=Pt(12), color=DARK_TEXT, align=PP_ALIGN.CENTER)
        sx += stat_bw + int(0.10*IN)

    t = slide.shapes.add_textbox(bx, by + int(1.32*IN), bw, h - hh - int(1.38*IN))
    tf = t.text_frame; tf.word_wrap = True
    lines = [
        ("Token Reduction via Mini-Tables", True, Pt(13.5), NAVY),
        ("Mini-tables (top-5 rows) cut context length by >70% vs. full tables:", False, Pt(13), DARK_TEXT),
        ("  Llama3: 1,783 → 363 tok (k=1)  ·  Mistral: 3,314 → 477 tok", False, Pt(13), DARK_TEXT),
        ("  Qwen2.5: 1,951 → 394 tok  →  less noise for downstream LLM", False, Pt(13), DARK_TEXT),
        ("", False, Pt(3), DARK_TEXT),
        ("Embedding Call Reduction", True, Pt(13.5), NAVY),
        ("At inference CRAFT executes only 3 query embeddings + 5,100 mini-table embeddings per query:", False, Pt(13), DARK_TEXT),
        ("  Stage 1: 1 query encoding + inverted-index lookup", False, Pt(13), DARK_TEXT),
        ("  Stage 2: 1 ANN search over pre-encoded row embeddings", False, Pt(13), DARK_TEXT),
        ("  Stage 3: 1 query embed + re-rank top-100 mini-tables", False, Pt(13), DARK_TEXT),
        ("Dense baseline requires 169,898 table embedding calls — CRAFT is 33× leaner.", True, Pt(13), GREEN_DARK),
        ("", False, Pt(3), DARK_TEXT),
        ("Adaptability", True, Pt(13.5), NAVY),
        ("Stage 3 re-ranker is plug-and-play: swap in text-embedding-3-large,\n"
         "gemini-embedding-001, or future models with zero architectural changes.", False, Pt(13), DARK_TEXT),
    ]
    first = True
    for txt, bold, sz, clr in lines:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(1)
        r = p.add_run(); r.text = txt
        r.font.size = sz; r.font.bold = bold; r.font.color.rgb = clr


# ── Section: Conclusion ────────────────────────────────────────────────────

def sec_conclusion(slide, x, y, w, h):
    hh, by = sec_hdr(slide, x, y, w, "  Conclusion", fill=NAVY)
    rect(slide, x, by, w, h - hh, fill=BLUE_LIGHT, border=BLUE, bw=Pt(0.8))
    bx = x + int(0.14*IN); bw = w - int(0.28*IN)

    t = slide.shapes.add_textbox(bx, by + int(0.10*IN), bw, h - hh - int(0.14*IN))
    tf = t.text_frame; tf.word_wrap = True
    lines = [
        ("CRAFT is a modular, training-free cascaded retrieval framework for "
         "Open-Domain Table QA that achieves state-of-the-art recall without "
         "any dataset-specific fine-tuning.", False, Pt(13.5), DARK_TEXT),
        ("", False, Pt(3), DARK_TEXT),
        ("Key Takeaways", True, Pt(14), NAVY),
        ("▸  New SOTA Recall@1 = 49.84 on NQ-Tables", False, Pt(13.5), DARK_TEXT),
        ("▸  Competitive zero-shot performance on OTT-QA", False, Pt(13.5), DARK_TEXT),
        ("▸  33× fewer embedding calls vs. dense baselines", False, Pt(13.5), DARK_TEXT),
        ("▸  Robust to paraphrasing  (only –0.04 Δ recall)", False, Pt(13.5), DARK_TEXT),
        ("▸  Plug-and-play: swap any Stage-3 embedding model", False, Pt(13.5), DARK_TEXT),
        ("", False, Pt(3), DARK_TEXT),
        ("Impact: Bridges fine-tuned & lightweight retrieval — "
         "scalable, interpretable, and adaptable.", True, Pt(13.5), BLUE),
    ]
    first = True
    for txt, bold, sz, clr in lines:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(1)
        r = p.add_run(); r.text = txt
        r.font.size = sz; r.font.bold = bold; r.font.color.rgb = clr


# ── Main ───────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = Emu(W)
    prs.slide_height = Emu(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full background
    rect(slide, 0, 0, W, H, fill=RGBColor(0xEE, 0xF2, 0xF7))

    build_header(slide)

    # ── Layout ─────────────────────────────────────────────────────────────
    # Body area starts at BODY_Y, total height BODY_H
    # 3 columns: COL_X[0..2], each COL_W wide
    # Row heights (must sum to ≤ BODY_H with row gaps)
    # We use 4 horizontal bands:
    #   Band A (tall): left-col rows stacked, center = overview+pipeline, right = NQ+OTT
    #   Band B: e2e (spans 2 cols) + ablation (1 col)
    #   Band C: robustness + efficiency + conclusion

    BAND_A_H = int(27.5 * IN)   # overview figure + pipeline
    BAND_B_H = int( 8.5 * IN)   # e2e results
    BAND_C_H = BODY_H - BAND_A_H - BAND_B_H - 2*GAP

    AY = BODY_Y
    BY = AY + BAND_A_H + GAP
    CY = BY + BAND_B_H + GAP

    # ── Left column (band A) ──────────────────────────────────────────────
    left_sections = [
        (int(10.2 * IN), sec_motivation),
        (int( 5.8 * IN), sec_why_fails),
        (int( 5.0 * IN), sec_key_insight),
        (BAND_A_H - int(10.2*IN) - int(5.8*IN) - int(5.0*IN) - 3*GAP, sec_contributions),
    ]
    oy = AY
    for sh, fn in left_sections:
        fn(slide, COL_X[0], oy, COL_W, sh)
        oy += sh + GAP

    # ── Center column (band A): overview figure (top) + pipeline (rest) ──
    FIG_H = int(11.0 * IN)
    PIP_H = BAND_A_H - FIG_H - GAP
    sec_overview_fig(slide, COL_X[1], AY, COL_W, FIG_H)
    sec_pipeline(slide, COL_X[1], AY + FIG_H + GAP, COL_W, PIP_H)

    # ── Right column (band A): NQ + OTT ───────────────────────────────────
    NQ_H  = int(14.2 * IN)
    OTT_H = BAND_A_H - NQ_H - GAP
    sec_nq(slide, COL_X[2], AY, COL_W, NQ_H)
    sec_ottqa(slide, COL_X[2], AY + NQ_H + GAP, COL_W, OTT_H)

    # ── Band B: e2e (cols 0+1) + ablation (col 2) ─────────────────────────
    e2e_w = COL_W * 2 + GAP
    sec_e2e(slide, COL_X[0], BY, e2e_w, BAND_B_H)
    sec_ablation(slide, COL_X[2], BY, COL_W, BAND_B_H)

    # ── Band C: robustness | efficiency | conclusion ───────────────────────
    sec_robustness(slide, COL_X[0], CY, COL_W, BAND_C_H)
    sec_efficiency(slide, COL_X[1], CY, COL_W, BAND_C_H)
    sec_conclusion(slide, COL_X[2], CY, COL_W, BAND_C_H)

    # ── Footer ────────────────────────────────────────────────────────────
    FY = H - int(0.30*IN)
    rect(slide, 0, FY, W, int(0.30*IN), fill=NAVY, bw=Pt(0))
    tb(slide, int(0.3*IN), FY + int(0.03*IN), int(22*IN), int(0.26*IN),
       "CRAFT  ·  Training-Free Cascaded Retrieval for Tabular QA  ·  ACL 2026",
       size=Pt(12), color=RGBColor(0xCC, 0xD9, 0xE8))
    tb(slide, int(14*IN), FY + int(0.03*IN), int(21.5*IN), int(0.26*IN),
       "arXiv: 2505.14984   ·   asing725@asu.edu   ·   bhandk@rpi.edu   ·   "
       "coral-lab-asu.github.io",
       size=Pt(12), color=GOLD, align=PP_ALIGN.RIGHT)

    out = "/mnt/data1/asing725/ACL/CRAFT/CRAFT_poster.pptx"
    prs.save(out)
    sw = prs.slide_width.inches; sh2 = prs.slide_height.inches
    nshapes = len(slide.shapes)
    npics = sum(1 for s in slide.shapes if s.shape_type == 13)
    print(f"Saved: {out}")
    print(f"Size: {sw:.1f}\" × {sh2:.1f}\"  |  shapes: {nshapes}  |  images: {npics}")


if __name__ == "__main__":
    build()
